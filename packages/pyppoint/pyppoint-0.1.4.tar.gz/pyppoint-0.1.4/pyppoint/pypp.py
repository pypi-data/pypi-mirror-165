import logging
import subprocess

import treefiles as tf
from pptx import Presentation as PPresentation
from pptx.util import Inches, Pt


class Slide:
    def __init__(self, slide):
        self.slide = slide

    def image(self, fname, left, top, width=2):
        self.slide.shapes.add_picture(
            fname, Inches(left), Inches(top), width=Inches(width)
        )

    def text(self, txt, x, y, cx, cy, bold=False, size=20):
        textbox = self.slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(cx), Inches(cy)
        )
        textframe = textbox.text_frame
        p = textframe.add_paragraph()
        p.text = txt
        p.font.bold = bold
        p.font.size = Pt(size)


class Presentation:
    # default size: width=Inches(10), height=Inches(7.5)
    def __init__(self):
        self.prs = PPresentation()
        self.fname = None

    def title(self, title, subtitle=None):
        first_slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        first_slide.shapes.title.text = title
        if subtitle is not None:
            first_slide.placeholders[1].text = subtitle

    def slide(self) -> Slide:
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        return Slide(slide)

    def save(self, fname):
        self.fname = tf.ensure_ext(fname, "pptx")
        self.prs.save(self.fname)

    def open(self):
        assert self.fname is not None
        subprocess.Popen(["libreoffice", self.fname])


log = logging.getLogger(__name__)
