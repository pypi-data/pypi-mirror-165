from logging import raiseExceptions
import numpy as np
from wx import dataview
import wx
import wx.propgrid as pg
from wx.core import VERTICAL, BoxSizer, Height, ListCtrl, StaticText, TextCtrl, Width
from wx.glcanvas import GLCanvas,GLContext
from wx.dataview import TreeListCtrl
from PIL import Image
from PIL.PngImagePlugin import PngInfo
from io import BytesIO
from OpenGL.GL import *
from OpenGL.GLUT import *
import math
import matplotlib.pyplot as plt
from matplotlib.widgets import Button as mplButton
from os import scandir,listdir
from os.path import exists,join,normpath
from pptx import Presentation
import threading

from .PyPalette import wolfpalette
from .wolfresults_2D import Wolfresults_2D,CHOICES_VIEW_2D
from .PyTranslate import _
from .PyWMS import getIGNFrance,getWalonmap
from .PyVertex import cloud_vertices,getIfromRGB
from .RatingCurve import SPWMIGaugingStations,SPWDCENNGaugingStations
from .wolf_array import WOLF_ARRAY_MB, SelectionData, WolfArray, WolfArrayMB
from .PyParams import Wolf_Param
from .BcManager import BcManager
from .PyVertexvectors import *
from .Results2DGPU import wolfres2DGPU
from .PyCrosssections import crosssections, profile
from .GraphNotebook import PlotNotebook
from .laz_viewer import myviewer,read_laz,clip_data_xyz

ID_SELECTCS = 1000
ID_SORTALONG = 1001
ID_LOCMINMAX = 1002

class imagetexture():
    
    parent=None

    name:str
    idtexture:int
    xmin:float
    xmax:float
    ymin:float
    ymax:float

    width:int
    height:int

    which:str
    category:str
    subcategory:str

    France:bool
    epsg:str

    def __init__(self,which:str,label:str,cat:str,subc:str,parent,xmin,xmax,ymin,ymax,width=1000,height=1000,France=False,epsg='31370') -> None:
        
        parent.canvas.SetCurrent(parent.context)
        self.parent=parent
        self.France=France
        self.epsg=epsg
        
        self.plotted=False

        self.xmin=xmin
        self.xmax=xmax
        self.ymin=ymin
        self.ymax=ymax
        self.idtexture=(GLuint*1)()
        try:
            glGenTextures(1,self.idtexture)
        except:
            raise NameError('Opengl glGenTextures -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')
        self.width=width
        self.height=height
        self.which=which.lower()
        self.category=cat#.upper()
        self.name=label
        self.subcategory=subc#.upper()
        self.oldview=[self.xmin,self.xmax,self.ymin,self.ymax,self.width,self.height]

        self.load()        

        pass

    def load(self):
        if self.width==-99999 or self.height==-99999:
            return

        if self.parent.canvas.SetCurrent(self.parent.context):
            mybytes:BytesIO   

            if self.France:
                mybytes=getIGNFrance(self.category,self.epsg,
                                        self.xmin,self.ymin,self.xmax,self.ymax,
                                        self.width,self.height,False)
            else:
                mybytes=getWalonmap(self.category+'/'+self.subcategory,
                                        self.xmin,self.ymin,self.xmax,self.ymax,
                                        self.width,self.height,False)
            image=Image.open(mybytes)  
            
            glBindTexture(GL_TEXTURE_2D, self.idtexture[0])
            if self.subcategory[:5]=='ORTHO':
                glTexImage2D(GL_TEXTURE_2D, 0, GL_RGBA, self.width,self.height, 0, GL_RGBA, GL_UNSIGNED_BYTE, image.tobytes())
            elif image.mode=='RGB':
                glTexImage2D(GL_TEXTURE_2D, 0, GL_RGBA, self.width,self.height, 0, GL_RGB, GL_UNSIGNED_BYTE, image.tobytes())
            else:
                glTexImage2D(GL_TEXTURE_2D, 0, GL_RGBA, self.width,self.height, 0, GL_RGBA, GL_UNSIGNED_BYTE, image.tobytes())
            glTexParameteri(GL_TEXTURE_2D, GL_TEXTURE_MAG_FILTER, GL_LINEAR)
            glTexParameteri(GL_TEXTURE_2D, GL_TEXTURE_MIN_FILTER, GL_LINEAR_MIPMAP_LINEAR)
            glGenerateMipmap(GL_TEXTURE_2D)
        else:
            raise NameError('Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')
    
    def reload(self):
        dx=self.parent.xmax-self.parent.xmin
        dy=self.parent.ymax-self.parent.ymin
        cx=self.parent.mousex
        cy=self.parent.mousey

        coeff=.5
        self.xmin=cx-dx*coeff
        self.xmax=cx+dx*coeff
        self.ymin=cy-dy*coeff
        self.ymax=cy+dy*coeff
        self.width=self.parent.canvaswidth*2*coeff
        self.height=self.parent.canvasheight*2*coeff

        self.newview=[self.xmin,self.xmax,self.ymin,self.ymax,self.width,self.height]
        if self.newview!=self.oldview:
            self.load()        
            self.oldview=self.newview

    def paint(self):
        
        glPolygonMode(GL_FRONT_AND_BACK,GL_FILL) 
        
        glColor4f(1.,1.,1.,1.)
        glEnable(GL_TEXTURE_2D)
        glEnable(GL_BLEND)
        glBlendFunc(GL_SRC_ALPHA, GL_ONE_MINUS_SRC_ALPHA)

        glBindTexture(GL_TEXTURE_2D, self.idtexture[0])

        glBegin(GL_QUADS)
        glTexCoord2f(0.0, 0.0)
        glVertex2f(self.xmin,self.ymax)
        glTexCoord2f(1.0, 0.0)
        glVertex2f(self.xmax,self.ymax)
        glTexCoord2f(1.0, 1.0)
        glVertex2f(self.xmax,self.ymin)
        glTexCoord2f(0.0, 1.0)
        glVertex2f(self.xmin,self.ymin)
        glEnd()
        glDisable(GL_TEXTURE_2D)
        glDisable(GL_BLEND)
        glPolygonMode(GL_FRONT_AND_BACK,GL_LINE)        
    
    def check_plot(self):
        self.plotted = True
            
    def uncheck_plot(self,unload=True):
        self.plotted = False
        
class WolfMapViewer(wx.Frame):
    '''
    Fenêtre de visualisation de données WOLF grâce aux WxWidgets
    '''
    
    TIMER_ID = 100 #délai d'attente avant action
    
    mybc : BcManager = None #Gestionnaire de CL
    myarrays:list #matrices ajoutées
    myvectors:list #zones vectorielles ajoutées
    myclouds:list #nuages de vertices
    myothers:list 
    mywmsback:list 
    mywmsfore:list 
    myres2D:list 

    canvas:GLCanvas         #canvas OpenGL
    context:GLContext       #context OpenGL
    mytooltip:Wolf_Param    #Objet WOLF permettant l'analyse de ce qui est sous la souris
    treelist:TreeListCtrl   #Gestion des éléments sous forme d'arbre
    selitem:StaticText
    leftbox:BoxSizer
    added:dict              # dictionnaire des éléments ajoutés
    
    active_vector:vector
    active_zone:zone
    active_zones:Zones
    active_array:WolfArray
    active_vertex:wolfvertex
    active_cs:crosssections
    

    def __init__(self, parent, title,w=500,h=500,treewidth=200):
        
        self.action = None #Action à entreprendre
        self.update_absolute_minmax=False #Force la MAJ de la palette
        self.copyfrom=None #aucun élément pointé par CTRL+C
        
        self.regular = True #Gestion de la taille de fenêtre d'affichage, y compris l'arbre de gestion
        self.sx=1                #facteur d'échelle selon X = largeur en pixels/largeur réelle
        self.sy=1                #facteur d'échelle selon Y = hauteur en pixels/hauteur réelle
        self.samescale = True    #force le même facteur d'échelle
        
        self.dynapar_dist = 1.

        #emprise initiale
        self.xmin=0.
        self.ymin=0.
        self.xmax=40.
        self.ymax=40.
        self.width = self.xmax-self.xmin   #largeur de la zone d'affichage en coordonnées réelles
        self.height = self.ymax-self.ymin  #hauteur de la zone d'affichage en coordonnées réelles
        self.canvaswidth=100
        self.canvasheight=100

        #position de la caméra
        self.mousex=self.width/2.
        self.mousey=self.height/2.

        self.bordersize = 0          #zone réservée au contour
        self.titlesize = 0           #zone réservée au titre
        self.treewidth = 200         #largeur de la zone d'arbre "treelist"

        self.backcolor = wx.Colour(255,255,255)  #couleur de fond
        self.mousedown=(0.,0.)                   #position initiale du bouton position bas
        self.mouseup=(0.,0.)                     #position initiale du bouton position haut
        self.oneclick=True   #détection d'un simple click ou d'un double-click
        self.move=False      #la souris est-elle en train de bouger?
        
        self.linked=False
        self.link_shareopsvect=True
        self.linkedList=None
        self.link_params=None
        
        self.forcemimic=True
        self.currently_readresults=False
        
        self.mylazdata=None
        
        self.treewidth = treewidth
        super(WolfMapViewer, self).__init__(parent, title = title,size = (w+self.treewidth,h))

        #Gestion des menus
        self.popupmenu = wx.Menu()
        for text in "one two three four five".split():
            item = self.popupmenu.Append(-1, text)
            self.Bind(wx.EVT_MENU, self.OnPopupItemSelected, item)

        self.menubar = wx.MenuBar()
        
        self.menuwolf2d=None

        self.filemenu = wx.Menu()
        openitem = self.filemenu.Append(wx.ID_OPEN,_('Open project'),_('Open project'))
        # saveproject = self.filemenu.Append(wx.ID_ANY,_('Save project'),_('save project'))
        self.filemenu.AppendSeparator()
        saveitem = self.filemenu.Append(wx.ID_SAVE,_('Save'),_('Save data to files'))
        saveasitem = self.filemenu.Append(wx.ID_SAVEAS,_('Save as...'),_('Save data to new files'))
        savecanvas = self.filemenu.Append(wx.ID_ANY,_('Save canvas...'),_('Save canvas to image file'))
        self.filemenu.AppendSeparator()
        exportgltf = self.filemenu.Append(wx.ID_ANY,_('Export to gltf...'),_('Save data to gltf files'))
        importgltf = self.filemenu.Append(wx.ID_ANY,_('Import from gltf...'),_('Import data from gltf files'))
        compareitem = self.filemenu.Append(wx.ID_ANY,_('Set gltf comparison'),_('Set gltf comparison'))
        updategltf = self.filemenu.Append(wx.ID_ANY,_('Update from gltf...'),_('Update data from gltf files'))
        self.filemenu.AppendSeparator()
        compareitem = self.filemenu.Append(wx.ID_ANY,_('Set comparison'),_('Set comparison'))
        multiview = self.filemenu.Append(wx.ID_ANY,_('Multiviewer'),_('Multiviewer'))
        self.filemenu.AppendSeparator()
        createarray = self.filemenu.Append(wx.ID_FILE6,_('Create array...'),_('New array (binary file - real)'))
        createvector = self.filemenu.Append(wx.ID_FILE7,_('Create vectors...'),_('New vectors'))
        createcloud = self.filemenu.Append(wx.ID_FILE8,_('Create cloud...'),_('New cloud'))
        self.filemenu.AppendSeparator()
        addarray = self.filemenu.Append(wx.ID_FILE1,_('Add array...'),_('Add array (binary file - real)'))
        addvector = self.filemenu.Append(wx.ID_FILE2,_('Add vectors...'),_('Add vectors'))
        addcloud = self.filemenu.Append(wx.ID_FILE3,_('Add cloud...'),_('Add cloud'))
        addprofiles = self.filemenu.Append(wx.ID_FILE4,_('Add cross sections...'),_('Add cross sections'))
        addres2D = self.filemenu.Append(wx.ID_ANY,_('Add Wolf2D results...'),_('Add Wolf 2D results'))
        self.filemenu.AppendSeparator()
        addscan = self.filemenu.Append(wx.ID_FILE5,_('Recursive scan...'),_('Add recursively'))

        self.toolsmenu = wx.Menu()
        self.select_cs = self.toolsmenu.Append(ID_SELECTCS,_("Select cross section"),_("Select cross section"), kind = wx.ITEM_CHECK)
        self.sortalong = self.toolsmenu.Append(ID_SORTALONG,_("Sort along..."),_("Sort cross sections along support vector"))
        self.menumanagebanks = self.toolsmenu.Append(wx.ID_ANY,_("Manage banks..."),_("Manage banks"))
        self.renamecs = self.toolsmenu.Append(wx.ID_ANY,_("Rename cross sections..."),_("Rename"))
        self.menutrianglecs = self.toolsmenu.Append(wx.ID_ANY,_("Triangulate cross sections..."),_("Triangulate"))
        self.menuviewerinterpcs=None
        self.menuinterpcs = None
        
        self.minmaxmenu = wx.Menu()
        self.locminmax = self.minmaxmenu.Append(ID_LOCMINMAX,_("Local minmax"),_("Adapt palette on current zoom"), kind = wx.ITEM_CHECK)

        self.lazmenu = wx.Menu()
        croplaz = self.lazmenu.Append(wx.ID_ANY,_('Clip LAZ data on current zoom'),_('Clip LAZ'))
        readlaz = self.lazmenu.Append(wx.ID_ANY,_('Read LAZ data from npz'),_('read LAZ'))
        viewlaz = self.lazmenu.Append(wx.ID_ANY,_('Create LAZ viewer'),_('LAZ Viewer'))
        bridgelaz = self.lazmenu.Append(wx.ID_ANY,_('Create cloud points from bridges'),_('LAZ Bridge'))
        buildinglaz = self.lazmenu.Append(wx.ID_ANY,_('Create cloud points from buildings'),_('LAZ Buildings'))

                
        #Gestion des outils --> Utile pour ManageActions
        self.tools={}
        curtool=self.tools[ID_SELECTCS]={}
        curtool['menu']=self.select_cs
        curtool['name']='Select nearest profile'
        
        self.active_vector = None
        self.active_zone = None
        self.active_zones=None
        self.active_vertex = None
        self.active_array = None
        self.active_cs = None
        self.active_res2d = None

        curtool=self.tools[ID_SORTALONG]={}
        curtool['menu']=self.sortalong
        curtool['name']='Sort along vector'

        curtool=self.tools[ID_LOCMINMAX]={}
        curtool['menu']=self.locminmax
        curtool['name']=None

        self.menubar.Append(self.filemenu,_('&File'))
        self.menubar.Append(self.toolsmenu,_('&Cross sections'))
        self.menubar.Append(self.lazmenu,_('&LAZ data'))
        self.menubar.Append(self.minmaxmenu,_('&Palette'))
        self.SetMenuBar(self.menubar)
        self.Bind(wx.EVT_MENU, self.OnMenubar)

        #Ajout du conteneur OpenGL
        self.canvas = GLCanvas(self)
        self.context = GLContext(self.canvas)
        self.mybackisloaded=False
        self.myfrontisloaded=False

        #ajout d'une liste en arbre des objets
        self.treelist = TreeListCtrl(self,style=wx.dataview.TL_CHECKBOX|wx.LC_EDIT_LABELS)
        self.selitem = StaticText(self)

        self.root = self.treelist.GetRootItem()
        self.treelist.AppendColumn(_('Objects to plot'))
        self.myitemsarray=self.treelist.AppendItem(self.root, _("Arrays"))
        self.myitemsvector=self.treelist.AppendItem(self.root, _("Vectors"))
        self.myitemscloud=self.treelist.AppendItem(self.root, _("Clouds"))
        self.myitemsres2d=self.treelist.AppendItem(self.root, _("Wolf2D"))
        self.myitemsothers=self.treelist.AppendItem(self.root, _("Others"))
        self.myitemswmsback=self.treelist.AppendItem(self.root, _("WMS-background"))
        self.myitemswmsfore=self.treelist.AppendItem(self.root, _("WMS-foreground"))

        self.added={}
        self.added['arrays']={}
        self.added['vectors']={}
        self.added['clouds']={}
        self.added['wolf2d']={}
        self.added['others']={}
        self.added['wms-background']={}
        self.added['wms-foreground']={}

        width,height=self.GetClientSize()
        self.bordersize = int((w-width+self.treewidth)/2)
        self.titlesize = h-height-self.bordersize
        self.SetSize(w+2*self.bordersize+self.treewidth,h+self.bordersize+self.titlesize)

        #dimensionnement et positionnement de la fenêtre OpenGL
        self.canvas.SetSize(width-self.treewidth,height)
        self.canvas.SetPosition((self.treewidth,0))

        self.setbounds()

        #dimensionnement et positionnement de l'arbre
        self.leftbox=BoxSizer(orient=wx.VERTICAL)
        self.leftbox.Add(self.treelist,1,wx.LEFT)
        self.leftbox.Add(self.selitem,0,wx.LEFT)
        self.treelist.SetSize(self.treewidth,height)
        #self.selitem.SetSize(self.treewidth,30)
        self.SetSizer(self.leftbox)
        
        #self.treelist.SetPosition((0,0))
        #self.selitem.SetPosition((0,height-30))
    
        #fenêtre ToolTip
        self.mytooltip = Wolf_Param(self,"Values",ontop=True,to_read=False,withbuttons=False)
        self.mytooltip.SetSize(300,400)
        self.mytooltip.Show(False)

        self.timer = wx.Timer(self, self.TIMER_ID)
        self.Bind(wx.EVT_TIMER, self.OnLMaintains)

        self.notebookcs = None
        self.notebookbanks = None
        self.myaxcs = None
        self.myfigcs = None
        
        self.InitUI()
    
    def menu_wolf2d(self):    
        
        if self.menuwolf2d is None:
            self.menuwolf2d = wx.Menu()
            self.menu2d_curentview = self.menuwolf2d.Append(wx.ID_ANY,_("Change current view"),_("Current view"))
            self.menu2d_lastres = self.menuwolf2d.Append(wx.ID_ANY,_("Read last result"),_("Current view"))
            self.menu2d_bc = self.menuwolf2d.Append(wx.ID_ANY,_("Manage boundary conditions..."),_("BC manager"))
            # self.menu2d_tft_bc = self.menuwolf2d.Append(wx.ID_ANY,_("Transfer boundary conditions..."),_("Transfer BC"))

            self.menubar.Append(self.menuwolf2d,_('Options 2D'))
    
    def triangulate_cs(self):
        
        msg=''
        if self.active_zones is None:
            msg+=_(' The active zones is None. Please activate the desired object !\n')
        if self.get_cross_sections() is None:
            msg+=_(' The is no cross section. Please load file!')

        if msg!='':
            dlg=wx.MessageBox(msg,'Required action')
            return

        dlg=wx.NumberEntryDialog(None,_('What is the desired size [cm] ?'),'ds','ds size',100,1,10000.)
        ret=dlg.ShowModal()
        if ret==wx.ID_CANCEL:
            dlg.Destroy()
            return
        
        ds=float(dlg.GetValue())/100.
        dlg.Destroy()
        
        self.myinterp = Interpolators(self.active_zones,self.get_cross_sections(),ds)
                            
        if self.menuviewerinterpcs is None:
            self.menuviewerinterpcs = self.toolsmenu.Append(wx.ID_ANY,_("New cloud Viewer..."),_("Cloud viewer Interpolate"))
        if self.menuinterpcs is None:
            self.menuinterpcs = self.toolsmenu.Append(wx.ID_ANY,_("Interpolate on active array..."),_("Interpolate"))

    def interpolate_cs(self):
        if self.active_array is not None and self.myinterp is not None:
            
            choices = ["nearest", "linear", "cubic"] 
            dlg = wx.SingleChoiceDialog(None, "Pick an interpolate method", "Choices", choices)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            method=dlg.GetStringSelection()
            dlg.Destroy()
            
            self.myinterp.interp_on_array(self.active_array,method)
    
    def save_canvasogl(self,fn='',mpl=True,ds=500.):
        
        self.Paint()
        
        if fn=='':
            dlg = wx.FileDialog(None,_('Choose file name'),wildcard='PNG (*.png)|*.png|JPG (*.jpg)|*.jpg',style=wx.FD_SAVE)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            fn=dlg.GetPath()
            dlg.Destroy()

        if self.canvas.SetCurrent(self.context):
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0, 0, self.canvaswidth, self.canvasheight, GL_RGBA, GL_UNSIGNED_BYTE)
            myimage:Image.Image
            myimage = Image.frombuffer("RGBA",(self.canvaswidth, self.canvasheight),data)
            myimage = myimage.transpose(1)

            metadata = PngInfo()
            metadata.add_text('xmin',str(self.xmin))
            metadata.add_text('ymin',str(self.ymin))
            metadata.add_text('xmax',str(self.xmax))
            metadata.add_text('ymax',str(self.ymax))

            if mpl:
                extent=(self.xmin,self.xmax,self.ymin,self.ymax)
                fig,ax=plt.subplots(1,1)
                ax.imshow(myimage,origin='upper',
                            extent=extent)

                x1=int((self.xmin//ds)*ds)
                x2=int((self.xmax//ds)*ds)
                y1=int((self.ymin//ds)*ds)
                y2=int((self.ymax//ds)*ds)
                x_label_list=np.linspace(x1,x2,int((x2-x1)/ds)+1,True)
                y_label_list=np.linspace(y1,y2,int((y2-y1)/ds)+1,True)
                ax.set_xticklabels(x_label_list)
                ax.set_yticklabels(y_label_list)
                plt.savefig(fn,dpi=150)
            else:
                myimage.save(fn,pnginfo=metadata)

            return fn
        else:
            raise NameError('Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')
    
    def reporting(self,dir=''):
        if dir=='':
            dlg = wx.DirDialog(None,"Choose directory to store reporting",style=wx.FD_SAVE)
            ret=dlg.ShowModal()
            
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            dir = dlg.GetPath()
            dlg.Destroy()
        
        myppt = Presentation(__file__)
        slide = myppt.slides.add_slide(0)

        
        for curzone in self.myzones:        
            for curvec in curzone.myvectors:
                curvec:vector
                if curvec.nbvertices>1:                
                    oldwidth=curvec.myprop.width
                    curvec.myprop.width=4
                    myname = curvec.myname

                    self.Activate_vector(curvec)
                    
                    if self.parentGUI.linked:
                        for curview in self.parentGUI.linkedList:
                            title = curview.GetTitle()
                            curview.zoomon_activevector()           
                            fn = path.join(dir,title + '_' + myname+'.png')
                            curview.save_canvasogl(fn)
                    else:
                        self.parentGUI.zoomon_activevector()           
                        fn = path.join(dir,myname+'.png')
                        self.parentGUI.save_canvasogl(fn)
                        
                        fn = path.join(dir,'palette_v_' + myname+'.png')
                        self.parentGUI.active_array.mypal.export_image(fn,'v')
                        fn = path.join(dir,'palette_h_' + myname+'.png')
                        self.parentGUI.active_array.mypal.export_image(fn,'h')
                        
                    curvec.myprop.width = oldwidth
        
    
    def InitUI(self):

        self.Bind(wx.EVT_SIZE, self.OnSize)
        self.Bind(wx.EVT_CLOSE , self.OnClose)

        self.canvas.Bind(wx.EVT_CONTEXT_MENU, self.OnShowPopup)
        self.canvas.Bind(wx.EVT_PAINT, self.OnPaint)
        self.treelist.Bind(wx.EVT_CHAR, self.OnHotKey)
        self.canvas.Bind(wx.EVT_CHAR, self.OnHotKey)
        self.canvas.Bind(wx.EVT_BUTTON, self.OnButton)
        self.canvas.Bind(wx.EVT_RIGHT_DCLICK, self.OnRDClick)
        self.canvas.Bind(wx.EVT_LEFT_DCLICK, self.OnLDClick)
        self.canvas.Bind(wx.EVT_LEFT_DOWN, self.OnLDown)
        self.canvas.Bind(wx.EVT_LEFT_UP, self.OnLUp)
        self.canvas.Bind(wx.EVT_MIDDLE_DOWN, self.OnLDown)
        self.canvas.Bind(wx.EVT_MIDDLE_UP, self.OnLUp)
        self.canvas.Bind(wx.EVT_RIGHT_DOWN, self.OnRightDown)
        self.canvas.Bind(wx.EVT_RIGHT_UP, self.OnRightUp)
        self.canvas.Bind(wx.EVT_MOTION, self.OnMotion)
        self.canvas.Bind(wx.EVT_LEAVE_WINDOW, self.OnLeave)
        self.canvas.Bind(wx.EVT_MOUSEWHEEL, self.OnButton)   
             
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CHECKED, self.OnCheckItem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_ACTIVATED,self.OnActivateTreeElem)
        # dispo dans wxpython 4.1 self.Bind(wx.EVT_GESTURE_ZOOM,self.OnZoomGesture)

        self.Centre()

        self.myarrays=[]
        self.myvectors=[]
        self.myclouds=[]
        self.myothers=[]
        self.mywmsback=[]
        self.mywmsfore=[]
        self.myres2D=[]

        self.Show(True)

    def OnSize(self, e):
        """
        Redimensionnement de la fenêtre
        """
        if self.regular:
            #retrouve la taille de la fenêtre
            width,height=self.GetClientSize()
            #enlève la barre d'arbre
            width-=self.treewidth
            #définit la taille de la fenêtre graphique OpenGL et sa position (à droite de l'arbre)
            self.canvas.SetSize(width,height)
            self.canvas.SetPosition((self.treewidth,0))
            #calcule les limites visibles sur base de la taille de la fenêtre et des coefficients sx sy
            self.setbounds()
            #fixe la taille de l'arbre (notamment la hauteur)
            #self.treelist.SetSize(self.treewidth,height)
            e.Skip()

    def ManageActions(self,id):
        
        curmenu = self.tools[id]['menu']
        
        if curmenu.IsCheckable():
            if not curmenu.IsChecked():
                curmenu.Check(False)
                self.action=None
                
                if id == ID_LOCMINMAX:
                    self.update_absolute_minmax = True             
            else:
                curmenu.Check()
                if not self.tools[id]['name'] is None:
                    self.action=self.tools[id]['name']
                    
        else:
            if id==ID_SORTALONG:
                #Tri le long d'un vecteur
                if not self.active_cs is None and not self.active_vector is None:
                    self.active_cs:crosssections
                    self.active_vector:vector
                    self.active_cs.sort_along(self.active_vector.asshapely_ls(),self.active_vector.myname,False)
                else:
                    msg=''
                    if self.active_cs is None:
                        msg+=_('Please select the active cross sections \n')
                    if self.active_vector is None:
                        msg+=_('Please select the active supprt vector')
                    mydiag = wx.MessageDialog(self,msg,_('Sort along'))
                    mydiag.ShowModal()
    
    def setbounds(self):
        '''
        Calcule les limites visibles de la fenêtrte graphique sur base des facteurs d'échelle courants
        '''
        self.updatescalefactors()
        #retrouve la taille de la fenêtre OpenGL
        width,height = self.canvas.GetSize()
        self.canvaswidth=width
        self.canvasheight=height

        #calcule la taille selon X et Y en coordonnées réelles
        width = width / self.sx
        height = height / self.sy
        #retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
        self.xmin = self.mousex - width/2.
        self.xmax = self.xmin+width
        self.ymin = self.mousey - height/2.
        self.ymax = self.ymin+height

        self.width = width
        self.height = height

        self.mousex=self.xmin+width/2.
        self.mousey=self.ymin+height/2.

        self.updatescalefactors()
        
        self.mybackisloaded=False
        self.myfrontisloaded=False
        
        self.Refresh()
        self.mimicme()

    def updatescalefactors(self):
        width, height = self.canvas.GetSize()

        self.sx=1
        self.sy=1
        if width>0:
            self.sx = float(width)/self.width
        if height>0:
            self.sy = float(height)/self.height

        self.sx=min(self.sx,self.sy)
        self.sy=self.sx

    def add_viewer_and_link(self):
        dlg=wx.TextEntryDialog(self,_('Enter a caption for the new window'))
        
        ret=dlg.ShowModal()
        
        if ret==wx.ID_CANCEL:
            dlg.Destroy()
            return
        
        newcap = dlg.GetValue()
        dlg.Destroy()
        newview=WolfMapViewer(None,newcap,w=600,h=600)
        
        if self.linkedList is None:
            self.linkedList = [self]
            
        self.linkedList.append(newview)
        
        for curview in self.linkedList:
            curview.linked=True
            curview.linkedList=self.linkedList
            curview.link_shareopsvect=False

    def add_grid(self):
        mygrid=Grid(1000.)
        self.add_object('vector',newobj=mygrid,ToCheck=False,id='Grid')

    def add_WMS(self):
        xmin=0
        xmax=0
        ymin=0
        ymax=0
        orthos={'IMAGERIE':{'1971':'ORTHO_1971','1994-2000':'ORTHO_1994_2000',
                        '2006-2007':'ORTHO_2006_2007',
                        '2009-2010':'ORTHO_2009_2010',
                        '2012-2013':'ORTHO_2012_2013',
                        '2015':'ORTHO_2015','2016':'ORTHO_2016','2017':'ORTHO_2017',
                        '2018':'ORTHO_2018','2019':'ORTHO_2019','2020':'ORTHO_2020',
                        '2021':'ORTHO_2021'}}
        for idx,(k,item) in enumerate(orthos.items()):
            for kdx,(m,subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                            newobj=imagetexture('PPNC',m,k,subitem,
                            self,xmin,xmax,ymin,ymax,-99999,1024),
                            ToCheck=False,id='PPNC '+m)
        self.add_object(which='wmsback',
                    newobj=imagetexture('PPNC','Orthos France','OI.OrthoimageCoverage.HR','',
                    self,xmin,xmax,ymin,ymax,-99999,1024,France=True,epsg='EPSG:27563'),
                    ToCheck=False,id='Orthos France')

        forelist={'EAU':{'Aqualim':'RES_LIMNI_DGARNE','Alea':'ALEA_INOND','Lidaxes':'LIDAXES'},
                    'LIMITES':{'Secteurs Statistiques':'LIMITES_QS_STATBEL'},
                    'INSPIRE':{'Limites administratives':'AU_wms'},
                    'PLAN_REGLEMENT':{'Plan Percellaire':'CADMAP_2021_PARCELLES'}}
        
        for idx,(k,item) in enumerate(forelist.items()):
            for kdx,(m,subitem) in enumerate(item.items()):
                self.add_object(which='wmsfore',
                            newobj=imagetexture('PPNC',m,k,subitem,
                            self,xmin,xmax,ymin,ymax,-99999,1024),
                            ToCheck=False,id=m)            

    def set_compare(self,ListArrays=None):
        #Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None,'Comparison',w=600,h=600)
        third = WolfMapViewer(None,'Difference',w=600,h=600)
        
        second.add_grid()
        third.add_grid()
        second.add_WMS()
        third.add_WMS()
        
        #Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list=[]
        list.append(first)
        list.append(second)
        list.append(third)
        
        #On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curlist in list:
            curlist.linked=True
            curlist.linkedList=list
        
        if ListArrays is not None:
            if len(ListArrays)==2:
                mnt = ListArrays[0]
                mns = ListArrays[1]
            else:
                return
        else:
            return

        mns:WolfArray
        mnt:WolfArray
        diff:WolfArray
        
        #Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = mns-mnt
        
        mns.copy_mask(mnt,True)
        diff.copy_mask(mnt,True)
        
        mnt.parentGUI = first
        mns.parentGUI = second
        diff.parentGUI = third

        mnt.myops.parentGUI = first
        mns.myops.parentGUI = second
        diff.myops.parentGUI = third
        
        path = os.path.dirname(__file__)
        fn=join(path,'models\\diff16.pal')
        diff.mypal.readfile(fn)
        diff.mypal.automatic=False
        diff.myops.palauto.SetValue(0)
        
        mnt.mypal.automatic=False
        mns.mypal.automatic=False
        mnt.myops.palauto.SetValue(0)
        mns.myops.palauto.SetValue(0)
        
        mnt.mypal.isopop(mnt.array,mnt.nbnotnull)
        mns.mypal.updatefrompalette(mnt.mypal)
        
        #Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array',newobj=mnt,ToCheck=True,id='source')
        second.add_object('array',newobj=mns,ToCheck=True,id='comp')
        third.add_object('array',newobj=diff,ToCheck=True,id='diff=comp-source')

        mnt.myops.myzones = mns.myops.myzones
        diff.myops.myzones = mns.myops.myzones
        
        first.active_array=mnt
        second.active_array=mns
        third.active_array=diff

    def set_compare_all(self,ListArrays=None):
        #Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None,'Comparison MNS',w=600,h=600)
        third = WolfMapViewer(None,'Difference MNS',w=600,h=600)
        if len(ListArrays)==3:
            fourth = WolfMapViewer(None,'Comparison MNT',w=600,h=600)
            fifth = WolfMapViewer(None,'Difference MNT',w=600,h=600)
        
        #Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list=[]
        list.append(first)
        list.append(second)
        list.append(third)
        if len(ListArrays)==3:
            list.append(fourth)
            list.append(fifth)

        for curview in list:
            if curview is not self:
                curview.add_grid()
                curview.add_grid()
        
        #On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curview in list:
            curview.linked=True
            curview.linkedList=list
        
        comp2=None
        if ListArrays is not None:
            if len(ListArrays)==2:
                src = ListArrays[0]
                comp1 = ListArrays[1]
            elif len(ListArrays)==3:
                src = ListArrays[0]
                comp1 = ListArrays[1]
                comp2 = ListArrays[2]
            else:
                return
        else:
            return

        src:WolfArray
        comp1:WolfArray
        diff1:WolfArray
        comp2:WolfArray
        diff2:WolfArray
        
        #Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff1 = comp1-src
        
        comp1.copy_mask(src,True)
        diff1.copy_mask(src,True)
        
        src.parentGUI = first
        comp1.parentGUI = second
        diff1.parentGUI = third

        src.myops.parentGUI = first
        comp1.myops.parentGUI = second
        diff1.myops.parentGUI = third
        
        src.mypal.automatic=False
        comp1.mypal.automatic=False
        src.myops.palauto.SetValue(0)
        comp1.myops.palauto.SetValue(0)
        
        src.mypal.isopop(src.array,src.nbnotnull)
        comp1.mypal.updatefrompalette(src.mypal)
        
        #Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array',newobj=src,ToCheck=True,id='source')
        second.add_object('array',newobj=comp1,ToCheck=True,id='comp')
        third.add_object('array',newobj=diff1,ToCheck=True,id='diff=comp-source')

        comp1.myops.myzones = src.myops.myzones
        diff1.myops.myzones = src.myops.myzones
        
        first.active_array=src
        second.active_array=comp1
        third.active_array=diff1

        if comp2 is not None:
            diff2 = comp2-src
            comp2.copy_mask(src,True)
            diff2.copy_mask(src,True)
            
            comp2.parentGUI = fourth
            diff2.parentGUI = fifth

            comp2.myops.parentGUI = fourth
            diff2.myops.parentGUI = fifth
            
            comp2.mypal.automatic=False
            comp2.myops.palauto.SetValue(0)
            
            comp2.mypal.updatefrompalette(src.mypal)
            
            #Ajout des matrices dans les fenêtres de visualisation
            fourth.add_object('array',newobj=comp2,ToCheck=True,id='comp2')
            fifth.add_object('array',newobj=diff2,ToCheck=True,id='diff2=comp2-source')

            comp2.myops.myzones = src.myops.myzones
            diff2.myops.myzones = src.myops.myzones
            
            fourth.active_array=comp2
            fifth.active_array=diff2

    def set_blender_sculpting(self):
        
        myframe = wx.Frame(None,title=_('Excavation and backfill'))
        sizergen = wx.BoxSizer(wx.VERTICAL)
        sizer1 = wx.BoxSizer(wx.HORIZONTAL)
        sizer2 = wx.BoxSizer(wx.HORIZONTAL)
        sizer3 = wx.BoxSizer(wx.HORIZONTAL)
        sizergen.Add(sizer1)
        sizergen.Add(sizer2)
        sizergen.Add(sizer3)
        
        labexc=wx.StaticText(myframe,label=_('Excavation : '))
        labback=wx.StaticText(myframe,label=_('Backfill   : '))
        labbal=wx.StaticText(myframe,label=_('Balance   : '))
        sizer1.Add(labexc)
        sizer2.Add(labback)
        sizer3.Add(labbal)
        
        font = wx.Font(18, wx.DECORATIVE, wx.NORMAL, wx.NORMAL)

        Exc=wx.StaticText(myframe,label=' [m³]')
        Back=wx.StaticText(myframe,label=' [m³]')
        Bal=wx.StaticText(myframe,label=' [m³]')
        
        labexc.SetFont(font)
        labback.SetFont(font)
        labbal.SetFont(font)
        Exc.SetFont(font)
        Back.SetFont(font)
        Bal.SetFont(font)
        
        sizer1.Add(Exc)
        sizer2.Add(Back)
        sizer3.Add(Bal)
        
        myframe.SetSizer(sizergen)
        myframe.Layout()
        myframe.Centre(wx.BOTH )
        myframe.Show()        
        
        if self.link_params is None:
            self.link_params={}        
        
        self.link_params['ExcavationBackfill'] = myframe    
        self.link_params['Excavation'] = Exc
        self.link_params['Backfill'] = Back
        self.link_params['Balance'] = Bal
        
        #Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None,'Sculpting',w=600,h=600)
        third = WolfMapViewer(None,'Difference',w=600,h=600)
        fourth = WolfMapViewer(None,'Gradient',w=600,h=600)
        fifth = WolfMapViewer(None,'Laplace',w=600,h=600)
        sixth = WolfMapViewer(None,'Unitary Mask',w=600,h=600)
        
        #Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list=[]
        list.append(first)
        list.append(second)
        list.append(third)
        list.append(fourth)
        list.append(fifth)
        list.append(sixth)

        for curlist in list:
            curlist.add_grid()
            curlist.add_WMS()        
        
        #On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curlist in list:
            curlist.linked=True
            curlist.linkedList=list
        
        source:WolfArray
        sourcenew:WolfArray
        diff:WolfArray
        grad:WolfArray
        lap:WolfArray
        unimask:WolfArray
        
        source = self.active_array
        sourcenew = WolfArray(mold=source)
        
        #Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = source-source
        grad = source.get_gradient_norm()
        lap = source.get_laplace()
        unimask = WolfArray(mold=diff)
        
        np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data!=0.)
        
        grad.copy_mask(source,True)
        lap.copy_mask(source,True)
        diff.copy_mask(source,True)
        unimask.copy_mask(source,True)
        
        sourcenew.parentGUI = second
        diff.parentGUI = third
        grad.parentGUI = fourth
        lap.parentGUI = fifth
        unimask.parentGUI = sixth

        sourcenew.myops.parentGUI = second
        diff.myops.parentGUI = third
        grad.myops.parentGUI = fourth
        lap.myops.parentGUI = fifth
        unimask.myops.parentGUI = sixth
        
        path = os.path.dirname(__file__)
        # fn=join(path,'models\\diff16.pal')
        # diff.mypal.readfile(fn)
        # diff.mypal.automatic=False
        # diff.myops.palauto.SetValue(0)

        # fn=join(path,'models\\diff3.pal')
        # unimask.mypal.readfile(fn)
        # unimask.mypal.automatic=False
        # unimask.myops.palauto.SetValue(0)
                        
        #Ajout des matrices dans les fenêtres de visualisation
        second.add_object('array',newobj=sourcenew,ToCheck=True,id='source_new')
        third.add_object('array',newobj=diff,ToCheck=True,id='diff=comp-source')
        fourth.add_object('array',newobj=grad,ToCheck=True,id='gradient')
        fifth.add_object('array',newobj=lap,ToCheck=True,id='laplace')
        sixth.add_object('array',newobj=unimask,ToCheck=True,id='unimask')

        sourcenew.myops.myzones = source.myops.myzones
        diff.myops.myzones = source.myops.myzones
        grad.myops.myzones = source.myops.myzones
        lap.myops.myzones = source.myops.myzones
        unimask.myops.myzones = source.myops.myzones
        
        second.active_array=sourcenew
        third.active_array=diff
        fourth.active_array=grad
        fifth.active_array=lap
        sixth.active_array=unimask
        
        self.mimicme()

    def update_blender_sculpting(self):
        
        if not self.linked:
            return
        if len(self.linkedList)!=6:
            return
        
        #Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self.linkedList[0]
        second = self.linkedList[1]
        third = self.linkedList[2]
        fourth = self.linkedList[3]
        fifth = self.linkedList[4]
        sixth = self.linkedList[5]
                
        source = first.active_array
        sourcenew = second.active_array
        diff = third.active_array
        grad = fourth.active_array
        lap = fifth.active_array
        unimask = sixth.active_array
        
        fn=''
        if self.link_params is not None:
            if 'gltf file' in self.link_params.keys():
                fn = self.link_params['gltf file']
                fnpos = self.link_params['gltf pos']
        
        if fn=='':    
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = self.link_params['gltf file']
                        fnpos = self.link_params['gltf pos']
                        break
        
        sourcenew.import_from_gltf(fn,fnpos)
                
        #Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff.array = (sourcenew-source).array
        grad.array = sourcenew.get_gradient_norm().array
        lap.array = sourcenew.get_laplace().array
        np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data!=0.)

        diff.copy_mask(sourcenew,True)
        lap.copy_mask(sourcenew,True)
        grad.copy_mask(sourcenew,True)
        unimask.copy_mask(sourcenew,True)
        
        first.Paint()
        second.Paint()
        third.Paint()
        fourth.Paint()
        fifth.Paint()
        sixth.Paint()
        
        Exc:wx.StaticText
        Back:wx.StaticText
        Bal:wx.StaticText
        if not 'ExcavationBackfill' in self.link_params.keys():
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'ExcavationBackfill' in curgui.link_params.keys():
                        myframe= curgui.link_params['ExcavationBackfill']
                        Exc = curgui.link_params['Excavation'] 
                        Back = curgui.link_params['Backfill'] 
                        Bal = curgui.link_params['Balance'] 
        else:
            myframe= self.link_params['ExcavationBackfill']
            Exc = self.link_params['Excavation'] 
            Back = self.link_params['Backfill'] 
            Bal = self.link_params['Balance'] 
        
        Exc.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array<0.]))+' [m³]')
        Back.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array>0.]))+' [m³]')
        Bal.SetLabel("{:.2f}".format(np.sum(diff.array))+' [m³]')
        
                        
    def zoomon_activevector(self,size=500.,forceupdate=True):
        
        curvec= self.active_vector
        if curvec.minx==-99999:
            curvec.find_minmax()
            
        bounds= [curvec.minx,curvec.maxx,curvec.miny,curvec.maxy]
        
        dx=bounds[1]-bounds[0]
        dy=bounds[3]-bounds[2]

        self.mousex = bounds[0]+dx/2.
        self.mousey = bounds[2]+dy/2.
        self.width = max(size,dx)
        self.height = max(size,dy)
        
        self.updatescalefactors()
        self.setbounds()
        self.mimicme()
        
        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def read_project(self,fn):
        myproject = Wolf_Param(None,filename=fn,toShow=False)

        mykeys=['cross_sections','vector','array']
        
        if 'which' in myproject.myparams.keys():
            which = myproject.myparams['which']['action']['value']
            if which == 'compare':
                ListCompare=[]
                if 'array' in myproject.myparams.keys():
                    for curid,curname in zip(myproject.myparams['array'].keys(),myproject.myparams['array'].values()):
                        ListCompare.append(WolfArray(normpath(curname['value'])))
                
                self.set_compare(ListCompare)
                return
        
        if 'cross_sections' in myproject.myparams.keys():
            mycs = crosssections(myproject.myparams['cross_sections']['mycs']['value'],
                                 format=myproject.myparams['cross_sections']['format']['value'],
                                dirlaz=myproject.myparams['cross_sections']['dirlaz']['value'])
            
            self.add_object('cross_sections',newobj=mycs,id='Cross Sections')

        if 'vector' in myproject.myparams.keys():
            myvec = Zones(myproject.myparams['vector']['river']['value'])
            self.add_object('vector',newobj=myvec,id='River')
            curvec:vector
            curvec=myvec.myzones[0].myvectors[0]
            mycs.sort_along(curvec.asshapely_ls(),curvec.myname,downfirst=False)

        if 'array' in myproject.myparams.keys():
            for curid,curname in zip(myproject.myparams['array'].keys(),myproject.myparams['array'].values()):
                curarray=WolfArray(curname['value'])
                self.add_object('array',newobj=curarray,id=curid)
                
        if 'palette' in myproject.myparams.keys():
            self.project_pal={}
            for curid,curname in zip(myproject.myparams['palette'].keys(),myproject.myparams['palette'].values()):               
                mypal = wolfpalette(None,'')
                mypal.readfile(curname['value'])
                mypal.automatic=False
                
                self.project_pal[curid]=mypal

        if 'palette-array' in myproject.myparams.keys():
            curarray:WolfArray
            if self.project_pal is not None:
                for curid,curname in zip(myproject.myparams['palette-array'].keys(),myproject.myparams['palette-array'].values()):               
                    if curname['value'] in self.project_pal.keys():                    
                        curarray=self.getobj(curid)
                        if curarray is not None:
                            mypal=self.project_pal[curname['value']]                       
                            curarray.mypal=mypal
                            curarray.myops.palauto.SetValue(0)
                            curarray.updatepalette(0)
                            curarray.delete_lists()

    def save_project(self,fn):
        myproject = Wolf_Param(None,toShow=False)

        mykeys=['cross_sections','vector','array']
        
        myproject.myparams['which']={}
        myproject.myparams['which']['action']={}
        myproject.myparams['which']['action']['value']

        mycs = self.get_cross_sections()
        if mycs is not None:
            myproject.myparams['cross_sections']={}
            myproject.myparams['cross_sections']['mycs']={}
            myproject.myparams['cross_sections']['mycs']['value']=mycs.filename

            myproject.myparams['vector']={}
            myproject.myparams['vector']['river']={}
            myproject.myparams['vector']['river']['value']=self.added['vectors'][0].filename

        if 'array' in myproject.myparams.key():
            for curid,curname in zip(myproject.myparams['array'].keys(),myproject.myparams['array'].values()):
                curarray=WolfArray(curname['value'])
                self.add_object('array',newobj=curarray,id=curid)

    def read_laz(self,fn=None):
        if fn is None:
            filternpz="npz (*.npz)|*.npz|all (*.*)|*.*"               
            dlg = wx.FileDialog(None,_('Choose LAS npz file'),wildcard=filternpz)
            ret=dlg.ShowModal()
            if ret != wx.ID_OK:
                return
            
            fn = dlg.GetPath()

        self.mylazdata=read_laz(fn)        

        if self.linked:
            if len(self.linkedList)>0:
                for curframe in self.linkedList:
                    curframe.mylazdata=self.mylazdata

    def managebanks(self):
        if self.notebookbanks is None:
            self.notebookbanks = PlotNotebook()
            self.mypagebanks  = self.notebookbanks.add(_("Manager banks interpolator"),"ManagerInterp")
                        
        msg=''
        if self.active_zones is None:
            msg+=_(' The active zones is None. Please activate the desired object !\n')
        if self.active_zone is None:
            msg+=_(' The active zone is None. Please activate the desired object !\n')
        if self.get_cross_sections() is None:
            msg+=_(' The is no cross section. Please load file!')

        if msg!='':
            dlg=wx.MessageBox(msg,'Required action')
            return
        
        self.mypagebanks.pointing(self,self.active_zones,self.active_zone,self.get_cross_sections())
        self.notebookbanks.Show()

    def _set_fn_fnpos_gltf(self):
            dlg=wx.FileDialog(None,_('Choose filename'),wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*',style=wx.FD_OPEN)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            fn=dlg.GetPath()
            dlg.Destroy()

            dlg=wx.FileDialog(None,_('Choose pos filename'),wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',style=wx.FD_OPEN)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            fnpos=dlg.GetPath()
            dlg.Destroy()
            
            if self.link_params is None:
                self.link_params={}
                
            self.link_params['gltf file'] = fn
            self.link_params['gltf pos'] = fnpos
            
            return fn

    def set_fn_fnpos_gltf(self):
            fn=''
            fnpos=''
            if self.linked:
                for curgui in self.linkedList:
                    if curgui.link_params is not None:
                        if 'gltf file' in curgui.link_params.keys():
                            fn = curgui.link_params['gltf file']
                            fnpos = curgui.link_params['gltf pos']
                            break
            elif self.link_params is None:
                self.link_params={}
                fn=self._set_fn_fnpos_gltf()

            if fn=='':
                self._set_fn_fnpos_gltf()
            
    def OnMenubar(self,event:wx.CommandEvent):
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())
        itemlabel=item.ItemLabel
        
        autoscale = True

        if id==wx.ID_OPEN:
            filterProject="proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"        
            file=wx.FileDialog(self,"Choose file", wildcard=filterProject)      
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                #récuparétaion du nom de fichier avec chemin d'accès
                filename =file.GetPath()
                file.Destroy()

            os.chdir(os.path.dirname(filename))                
            self.read_project(filename) 
        
        elif itemlabel==_("Change current view"):

            dlg = wx.SingleChoiceDialog(None, "Pick a view", "Choices", CHOICES_VIEW_2D)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            method=dlg.GetStringSelection()
            dlg.Destroy()
            
            self.active_res2d.set_currentview(method)
            
        elif itemlabel==_("Read last result"):
            
            self.currently_readresults=True
            
            for curmodel in self.added['wolf2d']:
                if self.added['wolf2d'][curmodel]['checked']:
                    locmodel:Wolfresults_2D
                    locmodel=self.added['wolf2d'][curmodel]['values']                                        
                    locmodel.read_oneresult()
                    locmodel.set_currentview(locmodel.currentview)

            self.currently_readresults=False

        elif itemlabel==_("Manage boundary conditions..."):
            if self.active_res2d is not None:
                
                self.active_res2d.myparams.editing_bc(self.added['wolf2d'])

        elif itemlabel==_("Manage banks..."):
            
            self.managebanks()
        
        elif itemlabel==_("Rename cross sections..."):
            
            dlg = wx.TextEntryDialog(None,_('Which starting point?'))
            ret=dlg.ShowModal()
            
            idxstart = dlg.GetValue()
            
            self.get_cross_sections().rename(int(idxstart))
            
        elif itemlabel==_("Triangulate cross sections..."):
            self.triangulate_cs()
        
        elif itemlabel==_("New cloud Viewer..."):
            if self.myinterp is not None:
                self.myinterp.viewer_interpolator()
            
        elif itemlabel==_("Interpolate on active array..."):
            if self.myinterp is not None:
                self.interpolate_cs()
            
        elif itemlabel==_('Save project'):
            filterProject="proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"        
            file=wx.FileDialog(self,"Name your file", wildcard=filterProject,style=wx.FD_SAVE)      
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                #récuparétaion du nom de fichier avec chemin d'accès
                filename =file.GetPath()
                file.Destroy()
                
            self.save_project(filename) 
        elif itemlabel == _('Read LAZ data from npz'):
            self.read_laz()

        elif itemlabel==_('Create cloud points from bridges'):
            if self.mylazdata is None:
                self.read_laz()
                
            mybridges = self.mylazdata[np.where(self.mylazdata[:,3]==10)] 
            mycloud = cloud_vertices()
            mycloud.init_from_nparray(mybridges)
            mycloud.myprop.style=2
            mycloud.myprop.color = getIfromRGB([255,102,102])
            mycloud.myprop.width=.5
            
            if self.linked:
                if len(self.linkedList)>0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud',newobj=mycloud,ToCheck=True,id='Bridges')                            
            else:
                self.add_object('cloud',newobj=mycloud,ToCheck=True,id='Bridges')
                   
        elif itemlabel==_('Create cloud points from buildings'):
            if self.mylazdata is None:
                self.read_laz()
                
            mybuildings = self.mylazdata[np.where(self.mylazdata[:,3]==1)] 
            mycloud = cloud_vertices()
            mycloud.init_from_nparray(mybuildings)
            mycloud.myprop.style=2
            mycloud.myprop.color = getIfromRGB([102,102,102])
            mycloud.myprop.width=.5
            if self.linked:
                if len(self.linkedList)>0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud',newobj=mycloud,ToCheck=True,id='Buildings')                            
            else:
                self.add_object('cloud',newobj=mycloud,ToCheck=True,id='Buildings')
                   
        elif itemlabel == _('Create LAZ viewer'):
            if self.mylazdata is None:
                self.read_laz()
            myviewer(self.mylazdata,0)     
            
        elif itemlabel == _('Clip LAZ data on current zoom'): 
            filternpz="npz (*.npz)|*.npz|all (*.*)|*.*"               
            dlg = wx.DirDialog(None,_('Choose XYZ Numpy LAZ Directory'))
            ret=dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            dlg = wx.FileDialog(None,_('Choose LAS npz out file'),wildcard=filternpz,style=wx.FD_SAVE)
            ret=dlg.ShowModal()
            if ret != wx.ID_OK:
                return
            
            fn = dlg.GetPath()
            
            self.mylazdir = dlg.GetPath()
            curbounds = [[self.xmin,self.xmin+self.width],[self.ymin,self.ymin+self.height]]
            clip_data_xyz(self.mylazdir,fn,curbounds)
            self.mydatalaz = read_laz(fn)
            
        elif itemlabel==_('Multiviewer'):
            dlg=wx.NumberEntryDialog(self,_("Additional viewers"),_("How many?"),_("How many additional viewers?"),1,0,5)
            ret=dlg.ShowModal()
            
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            nb = dlg.GetValue()
            dlg.Destroy()
            for i in range(nb):
                self.add_viewer_and_link()
                    
        elif itemlabel==_('Set Comparison'):
            ListArrays=[]
            filterProject="Array (*.bin)|*.bin|all (*.*)|*.*"        
            file=wx.FileDialog(self,"Choose Source file", wildcard=filterProject)      
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                #récuparétaion du nom de fichier avec chemin d'accès
                ListArrays.append(file.GetPath())
                file.Destroy()

            file=wx.FileDialog(self,_("Choose Comparison file"), wildcard=filterProject)      
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                #récuparétaion du nom de fichier avec chemin d'accès
                ListArrays.append(file.GetPath())
                file.Destroy()
            
            first=WolfArray(ListArrays[0],preload=False)
            sec=WolfArray(ListArrays[1],preload=False)
            first.read_txt_header()
            sec.read_txt_header()

            if first.nbx == sec.nbx and first.nby == sec.nby and first.origx == sec.origx and first.origy == sec.origy:
                self.set_compare(ListArrays)
                
        elif id==wx.ID_FILE1 :
            self.add_object(which='array',ToCheck=True)
            #self.mybc=BcManager(self,newarray.dx,newarray.dy,newarray.origx,newarray.origy,"Boundary conditions")
            #self.mybc.FindBorders(newarray.array)
        elif id==wx.ID_FILE2:
            self.add_object(which='vector',ToCheck=True)
        elif id==wx.ID_FILE3:
            self.add_object(which='cloud',ToCheck=True)
        elif id==wx.ID_FILE4:
            self.add_object(which='cross_sections',ToCheck=True)
        elif itemlabel==_('Add Wolf2D results...'):
            self.add_object(which='res2d',ToCheck=True)
            self.menu_wolf2d()
        elif id==wx.ID_FILE5:
            def addscandir(mydir):                
                for entry in scandir(mydir):
                    if entry.is_dir():
                        addscandir(entry)
                    elif entry.is_file():
                        if entry.name.endswith('.vec') or entry.name.endswith('.vecz'):
                            self.add_object(which='vector',
                                            filename=join(mydir,entry.name),
                                            ToCheck=True,
                                            id=join(mydir,entry.name))
                        elif entry.name.endswith('.bin'):
                            self.add_object(which='array',
                                            filename=join(mydir,entry.name),
                                            ToCheck=True,
                                            id=join(mydir,entry.name))

            mydialog=wx.DirDialog(self,_("Choose directory to scan"))
            if mydialog.ShowModal() == wx.ID_CANCEL:
                mydialog.Destroy()
                return
            else:
                #récuparétaion du nom de fichier avec chemin d'accès
                mydir =mydialog.GetPath()

            if exists(mydir):
                addscandir(mydir)
        
        elif id==wx.ID_FILE6:
            #Création d'une nouvelle matrice
            newarray = WolfArray(create=True,parentgui=self)
            self.add_object('array',newobj=newarray)
        elif id==wx.ID_FILE7:
            #Création de nouveaux vecteurs
            newzones = Zones(parent=self)
            self.add_object('vector',newobj=newzones)
        elif id==wx.ID_FILE8:
            #Création d'un nouveau nuage de point
            newcloud = cloud_vertices()
            self.add_object('cloud',newobj=newcloud)
        elif id in self.tools.keys():
            #gestion des actions
            self.ManageActions(id)
        elif id==wx.ID_SAVE:
            
            for curarray in self.added['arrays']: 
                if self.added['arrays'][curarray]['checked']:
                    obj:WolfArray
                    obj = self.added['arrays'][curarray]['values']
                    
                    if obj.filename=='':
                        filterArray="bin (*.bin)|*.bin|all (*.*)|*.*"
                        fdlg=wx.FileDialog(self,"Choose file", wildcard=filterArray,style=wx.FD_SAVE)
                        fdlg.ShowModal()
                        if fdlg.ShowModal() == wx.ID_OK:
                            obj.filename = fdlg.GetPath()
                            obj.write_all()      
            
            for curvect in self.added['vectors']:
                if self.added['vectors'][curvect]['checked']:
                    obj = self.added['vectors'][curvect]['values']
                    obj.saveas()      

        elif itemlabel=='Save canvas...':
            autoscale=False
            fn = self.save_canvasogl()
            self.save_linked_canvas(fn[:-4])
        elif itemlabel==_('Export to gltf...'):
            
            curarray:WolfArray
            curvec:vector
            
            msg=''
            if self.active_array is None:
                msg+=_('Active array is None\n')
            if self.active_vector is None:
                msg+=_('Active vector is None\n')
                
            if msg!='':
                msg+=_('\n')
                msg+=_('Retry !\n')
                wx.MessageBox(msg)
                return
            
            curarray = self.active_array
            curvec   = self.active_vector
            
            curvec.find_minmax()

            i1,j1 = curarray.get_ij_from_xy(curvec.minx,curvec.miny)
            x1,y1 = curarray.get_xy_from_ij(i1,j1)
            x1-=curarray.dx/2.
            y1-=curarray.dy/2.

            i2,j2 = curarray.get_ij_from_xy(curvec.maxx,curvec.maxy)
            x2,y2 = curarray.get_xy_from_ij(i2,j2)
            x2+=curarray.dx/2.
            y2+=curarray.dy/2.
            mybounds=[[x1,x2],[y1,y2]]
            
            dlg=wx.FileDialog(None,_('Choose filename'),wildcard='gltf2 (*.gltf)|*.gltf|All (*.*)|*.*',style=wx.FD_SAVE)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            fn=dlg.GetPath()
            dlg.Destroy()
            
            curarray.export_to_gltf(mybounds,fn)
            
        elif itemlabel==_('Import from gltf...'):

            curarray:WolfArray
            
            msg=''
            if self.active_array is None:
                msg+=_('Active array is None\n')
                
            if msg!='':
                msg+=_('\n')
                msg+=_('Retry !\n')
                wx.MessageBox(msg)
                return
            
            curarray = self.active_array
            
            dlg=wx.FileDialog(None,_('Choose filename'),wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*',style=wx.FD_OPEN)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            fn=dlg.GetPath()
            dlg.Destroy()

            dlg=wx.FileDialog(None,_('Choose pos filename'),wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',style=wx.FD_OPEN)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            fnpos=dlg.GetPath()
            dlg.Destroy()
            
            curarray.import_from_gltf(fn,fnpos)
                    
        elif itemlabel==_('Set gltf comparison'):
            
            self.set_blender_sculpting()
            autoscale=False
            
        elif itemlabel==_('Update from gltf...'):
            autoscale=False
            
            msg=''
            if self.active_array is None:
                msg+=_('Active array is None\n')
                
            if msg!='':
                msg+=_('\n')
                msg+=_('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_fn_fnpos_gltf()                                                   
            self.update_blender_sculpting()
                    
        elif id==wx.ID_SAVEAS:
            
            for curarray in self.added['arrays']: 
                if self.added['arrays'][curarray]['checked']:
                    obj:WolfArray
                    obj = self.added['arrays'][curarray]['values']
                    
                    filterArray="bin (*.bin)|*.bin|all (*.*)|*.*"
                    fdlg=wx.FileDialog(self,"Choose file name for Array : " + obj.idx, wildcard=filterArray,style=wx.FD_SAVE)
                    ret=fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        obj.filename = fdlg.GetPath()
                        obj.write_all()      
            
            for curvect in self.added['vectors']:
                if curvect=='grid':
                    pass
                elif self.added['vectors'][curvect]['checked']:
                    obj = self.added['vectors'][curvect]['values']
                    
                    filterArray="vec (*.vec)|*.vec|vec (*.vecz)|*.vecz|all (*.*)|*.*"
                    fdlg=wx.FileDialog(self,"Choose file name for Vector :"+ obj.idx, wildcard=filterArray,style=wx.FD_SAVE)
                    ret=fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        obj.saveas(fdlg.GetPath())      

        if len(self.myarrays)+len(self.myvectors)+len(self.myclouds)+len(self.myres2D)==2 and autoscale:
            #Trouve les bornzs si un seul élément est présent, sinon on conserve l'état du zoom
            self.Autoscale()

    def save_linked_canvas(self,fn):
        if self.linked:
            for idx,curel in enumerate(self.linkedList):
                curel.save_canvasogl(fn+'_'+str(idx)+'.png')

    def thread_update_blender(self):
        print("Update blender" )
        if self.canvas.SetCurrent(self.context):
            self.update_blender_sculpting()
            t = threading.Timer(10.0, self.thread_update_blender)
            t.start()

    def add_object(self,which:str='array',filename='',newobj=None,ToCheck=True,id=''):
        filterArray="bin (*.bin)|*.bin|all (*.*)|*.*"
        filterres2d="all (*.*)|*.*"
        filterVector="vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|all (*.*)|*.*"
        filterCloud="xyz (*.xyz)|*.xyz|text (*.txt)|*.txt|all (*.*)|*.*"
        filterCs="txt 2022 (*.txt)|*.txt|vecz WOLF (*.vecz)|*.vecz|WOLF (*.sxy)|*.sxy|text 2000 (*.txt)|*.txt|all (*.*)|*.*"

        if filename=='' and newobj is None:
            #ouverture d'une boîte de dialogue
            if which.lower()=='array':
                file=wx.FileDialog(self,"Choose file", wildcard=filterArray)
            elif which.lower()=='vector':
                file=wx.FileDialog(self,"Choose file", wildcard=filterVector)
            elif which.lower()=='cloud':
                file=wx.FileDialog(self,"Choose file", wildcard=filterCloud)
            elif which.lower()=='cross_sections':
                file=wx.FileDialog(self,"Choose file", wildcard=filterCs)
            elif which.lower()=='other':
                file=wx.FileDialog(self,"Choose file", wildcard=filterCloud)
            elif which.lower()=='res2d':
                file=wx.FileDialog(self,"Choose file", wildcard=filterres2d)
            elif which.lower()=='wmsback':
                file=wx.FileDialog(self,"Choose file", wildcard=filterCloud)
            elif which.lower()=='wmsfore':
                file=wx.FileDialog(self,"Choose file", wildcard=filterCloud)

            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                #récuparétaion du nom de fichier avec chemin d'accès
                filename =file.GetPath()
                curfilter = file.GetFilterIndex()
                file.Destroy()

        if filename !='':
            if(not(os.path.exists(filename))):
                print("Warning : the following file is not present here : " + filename)
                return 

        if which.lower()=='array':
            curdict=self.added['arrays']
            curtree=self.myitemsarray

            if newobj is None :
                
                testobj=WolfArray()
                testobj.filename = filename
                testobj.read_txt_header()
                
                if testobj.wolftype in WOLF_ARRAY_MB:
                    newobj = WolfArrayMB(filename,parentgui=self)
                else:
                    newobj = WolfArray(filename,parentgui=self)

            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            newobj.parentgui = self
            self.active_array = newobj
        
        elif which.lower()=='res2d':
            curdict=self.added['wolf2d']
            curtree=self.myitemsres2d

            if newobj is None :
                
                newobj=Wolfresults_2D(filename)
                newobj.read_param_simul()
                
            newobj.updatepalette()
            self.myres2D.append(newobj)
            newobj.parentgui = self
            self.active_res2d = newobj
            
        elif which.lower()=='vector':
            curdict=self.added['vectors']
            curtree=self.myitemsvector
            if newobj is None:
                newobj = Zones(filename)
            self.myvectors.append(newobj)

        elif which.lower()=='cross_sections':
            curdict=self.added['vectors']
            curtree=self.myitemsvector
                        
            if newobj is None:
                dlg=wx.DirDialog(None,'If exist, where are the LAZ data?')
                ret=dlg.ShowModal()
                if ret==wx.ID_OK:
                    dirlaz=dlg.GetPath()
                else:
                    dirlaz=''

                if curfilter==0: #txt 2022
                    newobj = crosssections(filename,format='2022',dirlaz=dirlaz)
                if curfilter==1: #vecz
                    newobj = crosssections(filename,format='vecz',dirlaz=dirlaz)
                elif curfilter==2: #sxy
                    newobj = crosssections(filename,format='sxy',dirlaz=dirlaz)
                else: #txt 2000
                    newobj = crosssections(filename,format='2000',dirlaz=dirlaz)
            
            self.myvectors.append(newobj)
        
        elif which.lower()=='cloud':
            curdict=self.added['clouds']
            curtree=self.myitemscloud
            if newobj is None:
                newobj = cloud_vertices(filename)
            self.myclouds.append(newobj)

        elif which.lower()=='other':
            if not newobj is None:
                curdict=self.added['others']
                curtree=self.myitemsothers
                self.myothers.append(newobj)

        elif which.lower()=='wmsback':
            if not newobj is None:
                curdict=self.added['wms-background']
                curtree=self.myitemswmsback
                self.mywmsback.append(newobj)

        elif which.lower()=='wmsfore':
            if not newobj is None:
                curdict=self.added['wms-foreground']
                curtree=self.myitemswmsfore
                self.mywmsfore.append(newobj)

        if id=='':
            dlg = wx.TextEntryDialog(self, 'ID ? (case insensitive)','Choose an identifier',_('NewObject'))
            dlg.SetValue('')
            if len(curdict)==0:
                if dlg.ShowModal() == wx.ID_OK:
                    id=dlg.GetValue()
            else:
                id=list(curdict.keys())[0]
                while id.lower() in curdict:
                    if dlg.ShowModal() == wx.ID_OK:
                        id=dlg.GetValue()
            dlg.Destroy()
            
        if id.lower() in curdict:
            endid='_'
            while id.lower() in curdict:
                id+=endid

        newobj.idx=id.lower()
            
        myitem=self.treelist.AppendItem(curtree,id,data=newobj)
        
        if ToCheck:
            self.treelist.CheckItem(myitem)
            self.treelist.CheckItem(self.treelist.GetItemParent(myitem))
            
            newobj.check_plot()

        loc=curdict[id.lower()]={}
        if filename!='':
            loc['listname']=filename.lower()
            loc['filename']=filename.lower()
        loc['values']=newobj
        loc['listitem']=myitem
        loc['checked']=ToCheck
        
        if type(newobj) == crosssections:
            self.myclouds.append(newobj.cloud)
            newobj.cloud.idx = newobj.idx
            newobj.cloud.myprop.filled=True
            myitem=self.treelist.AppendItem(self.myitemscloud,id,data=newobj.cloud)
                        
            curdict=self.added['clouds']
            loc=curdict[id.lower()]={}
            loc['values']=newobj.cloud
            loc['listitem']=myitem
            loc['checked']=False  
        elif type(newobj) == WolfArray:
            mycs = self.get_cross_sections()
            if mycs is not None:
                mycs.linked_arrays.append(newobj)
                mycs.linked_labels.append(id)                

    def getobj(self,id:str):
        for curdict in self.added:
            if id.lower() in self.added[curdict].keys():
                try:
                    return self.added[curdict][id]['values']
                except:
                    return None

    def OnShowPopup(self, event):
        pos = event.GetPosition()
        if pos==(-1,-1):
            width, height = self.GetSize()
            pos = (width/2.,height/2.)
        # else:
        #     pos = pos - self.GetPosition()
        self.PopupMenu(self.popupmenu, pos)

    def OnPopupItemSelected(self, event):
        item = self.popupmenu.FindItemById(event.GetId())
        text = item.GetText()
        wx.MessageBox("You selected item '%s'" % text)

    def OnClose(self,event):
        if self.linked:
            if self.linkedList is not None:
                if self in self.linkedList:
                    id=self.linkedList.index(self)
                    self.linkedList.pop(id)
            
        self.Destroy()

    def OnActivatedItem(self,e):
        myitem=e.GetItem()
        nameitem=self.treelist.GetItemText(myitem).lower()
        self.selitem.SetLabel(nameitem)
        pass

    def OnCheckItem(self,event):

        myitem=event.GetItem()
        myparent=self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)
        nameparent=self.treelist.GetItemText(myparent).lower()
        nameitem=self.treelist.GetItemText(myitem).lower()
        
        if nameparent != '':
            curobj=self.getobj(nameitem)
            self.added[nameparent][nameitem]['checked']=bool(check)
            
            if bool(check):
                curobj.check_plot()
            else:
                if type(curobj) in [WolfArray,WolfArrayMB]:
                    curobj.uncheck_plot(True)
                else:
                    curobj.uncheck_plot()
                
            if nameparent=='vectors' or nameparent=='cross_sections':
                curobj.showstructure(self)
            if curobj.idx=='grid' and check:
                dlg = wx.TextEntryDialog(self, 'Size of the Grid ? (float)','Choose an size')
                dlg.SetValue('1000.')
                size=1000.
                if dlg.ShowModal() == wx.ID_OK:
                    size=float(dlg.GetValue())
                curobj.creategrid(size,self.xmin,self.ymin,self.xmax,self.ymax)
                pass

    def getXY(self,pospix):

        width,height = self.canvas.GetSize()
        X=float(pospix[0])/self.sx+self.xmin
        Y=float(height - pospix[1])/self.sy+self.ymin
        return X,Y

    def OnZoomGesture(self,e):
        pass

    def OnLeave(self,e):
        if  e.ControlDown():
            self.mytooltip.Show(False)

    def get_cross_sections(self):
        for curvect in self.added['vectors']:
            if self.added['vectors'][curvect]['checked']:
                obj = self.added['vectors'][curvect]['values']
                if type(obj) is crosssections:
                    return obj
                    
        return None
        
    def OnRightDown(self,e:wx.MouseEvent):
        pos=e.GetPosition()
        x,y =self.getXY(pos)
        
        alt = e.AltDown()
        ctrl = e.ControlDown()
        
        if self.action is None:
            self.rightdown=(x,y)
            #self.OnShowPopup(e)            
        elif self.action=='Select nearest profile':
            
            if self.notebookcs is None:
                self.notebookcs = PlotNotebook()
                self.myfigcs  = self.notebookcs.add(_("Cross section"),"CS")
            else:
                try:
                    self.notebookcs.Show()
                except:
                    self.notebookcs = PlotNotebook()
                    self.myfigcs  = self.notebookcs.add(_("Cross section"),"CS")
                            
            curprofile:profile
            
            for curvect in self.added['vectors']:
                if self.added['vectors'][curvect]['checked']:
                    obj = self.added['vectors'][curvect]['values']
                    if type(obj) is crosssections:
                        if self.myfigcs.mycs is not None:
                            self.myfigcs.mycs.myprop.width=1
                            self.myfigcs.mycs.myprop.color=0
                        
                        curprofile = obj.select_profile(x,y)
                        self.myfigcs.set_cs(curprofile)
                        self.myfigcs.plot_cs()
                        
                        curprofile.myprop.width=2
                        curprofile.myprop.color=getIfromRGB([255,0,0])
                        
                        self.Paint()
        elif self.action.find('select active vector')>-1:
            
            inside = self.action.find('inside')>-1
            onlyonezone = self.action.find('2')>-1
            
            if onlyonezone:
                self.active_zone.select_vectors_from_point(x,y,inside)    
                self.active_vector = self.active_zone.get_selected_vectors()            

                if self.active_vector is not None:
                    self.active_zone.parent.Activate_vector(self.active_vector)
                    self.active_zone.active_vector = self.active_vector
                    self.active_zones.active_zone = self.active_vector.parentzone
            else:
                self.active_zones.select_vectors_from_point(x,y,inside)
                self.active_vector = self.active_zones.get_selected_vectors()            
            
            if self.active_vector is not None:
                self.active_zones.Activate_vector(self.active_vector)
                self.active_zone =self.active_vector.parentzone
                self.active_zones.expand_tree(self.active_zone)

        elif self.action=='select node by node':
            currray:WolfArray
            curarray=self.active_array
            curarray.mngselection.add_node_to_selection(x,y)
            curarray.mngselection.update_nb_nodes_sections()
            self.Paint()
        elif self.action=='select by tmp vector' or self.action=='select by vector':
            self.active_vector.add_vertex(wolfvertex(x,y)) 
        elif self.action=='capture vertices':
            self.active_vector.add_vertex(wolfvertex(x,y))  
        elif self.action=='dynamic parallel':
            
            if ctrl:
                if self.active_array is not None:
                    z = self.active_array.get_value(x,y)
                    self.active_vector.myvertices[-1].z=z
            
            self.active_vector.add_vertex(wolfvertex(x,y))                
            self.active_zone.parallel_active(self.dynapar_dist)
        elif self.action=='modify vertices':
            if self.active_vertex is None:
                self.active_vertex = self.active_vector.find_nearest_vert(x,y)
            else:
                self.active_vertex=None
        elif self.action=='insert vertices':
            if self.active_vertex is None:
                self.active_vertex = self.active_vector.insert_nearest_vert(x,y)
            else:
                self.active_vertex=None
        else:
            self.rightdown=(x,y)
            #self.OnShowPopup(e)

    def OnRightUp(self,e):
        pos=e.GetPosition()
        x,y =self.getXY(pos)

        try:
            minx=min(self.rightdown[0],x)
            miny=min(self.rightdown[1],y)
            maxx=max(self.rightdown[0],x)
            maxy=max(self.rightdown[1],y)

            if minx!=maxx and maxy!=miny :
                self.mybc.ray_tracing_numpy([[minx,miny],[maxx,miny],[maxx,maxy],[minx,maxy]],'X')
                self.mybc.ray_tracing_numpy([[minx,miny],[maxx,miny],[maxx,maxy],[minx,maxy]],'Y')
            else:
                self.mybc.query_kdtree((x,y))
        
            self.mybc.Populate()
            self.Refresh()

        except:
            pass

    def OnButton(self, e:wx.MouseEvent):
        d=e.GetWheelDelta()
        r=e.GetWheelRotation()
        a=e.GetWheelAxis()
        
        altdown = e.AltDown()
        ctrldown = e.ControlDown()

        if self.action=='dynamic parallel' and altdown and not ctrldown:
            self.dynapar_dist *= (1- .1 * (r/max(d,1))) 
            self.dynapar_dist = max(self.dynapar_dist,.01)
            
            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return
        elif self.action=='dynamic parallel' and altdown and ctrldown:
            dlg=wx.NumberEntryDialog(None,_('What is the desired size [cm] ?'),'ds','ds size',int(self.dynapar_dist*100.),1,100000.)
            ret=dlg.ShowModal()
            if ret==wx.ID_CANCEL:
                dlg.Destroy()
                return
            
            self.dynapar_dist=float(dlg.GetValue())/100.
            self.dynapar_dist = max(self.dynapar_dist,.01)
            dlg.Destroy()

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return
        
        self.width=self.width * (1- .1 * (r/max(d,1)))
        self.height=self.height * (1- .1 * (r/max(d,1)))

        self.setbounds()

    def OnRDClick(self,e):
        self.endactions()

    def OnLDClick(self,e):
        self.timer.Stop()
        pos=e.GetPosition()
        x,y =self.getXY(pos)
        self.mousex=self.mousedown[0]
        self.mousey=self.mousedown[1]
        self.mousedown=(0.,0.)
        self.oneclick=False
        self.setbounds()

    def OnLDown(self,e):
        if not self.move:
            pos=e.GetPosition()
            x,y =self.getXY(pos)
            self.mousedown=(x,y)
            self.mouseup=self.mousedown
        
            self.timer.Start(1000)

            e.Skip()
        else:
            self.move=False
            self.timer.Stop()
    
    def OnActivateTreeElem(self,e):
        curzones:Zones
        curzone:zone
        curvect:vector

        myitem=e.GetItem()
        myparent=self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)
        myobj=self.treelist.GetItemData(myitem)
        
        nameparent=self.treelist.GetItemText(myparent).lower()
        nameitem=self.treelist.GetItemText(myitem).lower()

        if type(myobj)==Zones:
            myobj.showstructure(self)
        elif type(myobj)==WolfArray:
            myobj.myops.Show()
            self.active_array=myobj
        elif type(myobj)==cloud_vertices:
            myobj.myprop.show()
        elif type(myobj)==crosssections:
            myobj.showstructure(self)
            self.active_cs=myobj
        elif type(myobj)==Wolfresults_2D:
            self.active_res2d=myobj
    def OnLMaintains(self,e):
        if not self.move:
            self.move=True
            print('move on!')
            self.timer.Start(10)
            e.Skip()
        else:
            dx=self.mouseup[0]-self.mousedown[0]
            dy=self.mouseup[1]-self.mousedown[1]

            self.mousex -= dx
            self.mousey -= dy
            self.mouseup = (self.mouseup[0]-dx,self.mouseup[1]-dy)
            
            self.mousedown=self.mouseup
            
            self.Refresh()
            e.Skip()

    def OnLUp(self,e):
        if not self.move:
            self.timer.Stop()
            pos=e.GetPosition()
            x,y =self.getXY(pos)
            self.mouseup = (x,y)

            if self.oneclick:
                self.mousex -= x-self.mousedown[0]
                self.mousey -= y-self.mousedown[1]
                self.setbounds()
            else:
                self.oneclick=True
        
        e.Skip()

    def OnMotion(self, e:wx.MouseEvent):

        #Déplacement de la souris sur le canvas OpenGL
        posframe = self.GetPosition()
        pos=e.GetPosition()
        x,y =self.getXY(pos)
        altdown = e.AltDown()
        
        if self.move:
            self.mouseup= (x,y)

        if self.action=='select by tmp vector' or self.action=='select by vector' or self.action=='capture vertices' or self.action=='dynamic parallel':
            if self.active_vector.nbvertices>0:
                self.active_vector.myvertices[-1].x=x
                self.active_vector.myvertices[-1].y=y

        if self.action=='modify vertices' or self.action=='insert vertices':
            if self.active_vertex is not None:
                if altdown:
                    ox=self.active_vector.myvertices[0].x
                    oy=self.active_vector.myvertices[0].y
                    
                    dirx = self.active_vector.myvertices[-1].x-ox
                    diry = self.active_vector.myvertices[-1].y-oy
                    normdir = np.sqrt(dirx**2.+diry**2.)
                    
                    dirx/=normdir
                    diry/=normdir
                    
                    vecx = x-ox
                    vecy = y-oy
                    
                    norm = np.sqrt(vecx**2.+vecy**2.)

                    self.active_vertex.x=ox + np.inner([dirx,diry],[vecx,vecy])* dirx
                    self.active_vertex.y=oy + np.inner([dirx,diry],[vecx,vecy])* diry
                    
                else:
                    self.active_vertex.x=x
                    self.active_vertex.y=y

        if self.action=='dynamic parallel':
            self.active_zone.parallel_active(self.dynapar_dist)
        
        self.mytooltip.myparams.clear()

        curgroup = 'Position'
        self.mytooltip.myparams[curgroup]={}

        curpar = 'Pixel'
        self.mytooltip.myparams[curgroup][curpar]={}
        self.mytooltip.myparams[curgroup][curpar]['name'] ='Pixel'
        self.mytooltip.myparams[curgroup][curpar]['value'] ='('+str(pos[0])+' ; '+str(pos[1])+')'
        self.mytooltip.myparams[curgroup][curpar]['type'] ='String'
        self.mytooltip.myparams[curgroup][curpar]['comment'] =''

        curpar = 'x'
        self.mytooltip.myparams[curgroup][curpar]={}
        self.mytooltip.myparams[curgroup][curpar]['name'] ='Coordinate X'
        self.mytooltip.myparams[curgroup][curpar]['value'] =x
        self.mytooltip.myparams[curgroup][curpar]['type'] ='Float'
        self.mytooltip.myparams[curgroup][curpar]['comment'] =''

        curpar = 'y'
        self.mytooltip.myparams[curgroup][curpar]={}
        self.mytooltip.myparams[curgroup][curpar]['name'] ='Coordinate Y'
        self.mytooltip.myparams[curgroup][curpar]['value'] =y
        self.mytooltip.myparams[curgroup][curpar]['type'] ='Float'
        self.mytooltip.myparams[curgroup][curpar]['comment'] =''

        for locarray in self.myres2D:
            curgroup = locarray.idx
            if self.added['wolf2d'][curgroup]['checked']:
                self.mytooltip.myparams[curgroup]={}

                try:
                    val = locarray.get_value(x,y,True)
                    i,j,curbloc = locarray.get_blockij_from_xy(x,y,False)
                    
                    if i!='-':
                        curpar = 'Indices'
                        self.mytooltip.myparams[curgroup][curpar]={}
                        self.mytooltip.myparams[curgroup][curpar]['name'] ='Indice (i;j;bloc)'
                        self.mytooltip.myparams[curgroup][curpar]['value'] ='('+str(i+1)+';'+str(j+1)+';'+str(curbloc)+')'
                        self.mytooltip.myparams[curgroup][curpar]['type'] ='String'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] =''

                        curpar = 'Value'
                        self.mytooltip.myparams[curgroup][curpar]={}
                        self.mytooltip.myparams[curgroup][curpar]['name'] ='Value'
                        self.mytooltip.myparams[curgroup][curpar]['value'] =float(val)
                        self.mytooltip.myparams[curgroup][curpar]['type'] ='Float'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] =''
                except:
                    pass

        for locarray in self.myres2D:
            curgroup = locarray.idx
            if self.added['wolf2d'][curgroup]['checked']:
                self.mytooltip.myparams[curgroup]={}

                try:
                    val = locarray.get_value(x,y,True)
                    i,j,curbloc = locarray.get_blockij_from_xy(x,y,False)
                    
                    if i!='-':
                        curpar = 'Indices'
                        self.mytooltip.myparams[curgroup][curpar]={}
                        self.mytooltip.myparams[curgroup][curpar]['name'] ='Indice (i;j;bloc)'
                        self.mytooltip.myparams[curgroup][curpar]['value'] ='('+str(i+1)+';'+str(j+1)+';'+str(curbloc)+')'
                        self.mytooltip.myparams[curgroup][curpar]['type'] ='String'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] =''

                        curpar = 'Value'
                        self.mytooltip.myparams[curgroup][curpar]={}
                        self.mytooltip.myparams[curgroup][curpar]['name'] ='Value'
                        self.mytooltip.myparams[curgroup][curpar]['value'] =float(val)
                        self.mytooltip.myparams[curgroup][curpar]['type'] ='Float'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] =''
                except:
                    pass

        for locarray in self.myarrays:
            curgroup = locarray.idx
            if self.added['arrays'][curgroup]['checked']:
                self.mytooltip.myparams[curgroup]={}

                try:
                    val = locarray.get_value(x,y)
                    if locarray.wolftype in WOLF_ARRAY_MB:
                        i,j,curbloc = locarray.get_blockij_from_xy(x,y)
                        curpar = 'Indices'
                        self.mytooltip.myparams[curgroup][curpar]={}
                        self.mytooltip.myparams[curgroup][curpar]['name'] ='Indice (i;j;bloc)'
                        self.mytooltip.myparams[curgroup][curpar]['value'] ='('+str(i+1)+';'+str(j+1)+';'+str(curbloc)+')'
                        self.mytooltip.myparams[curgroup][curpar]['type'] ='String'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] =''
                    else:
                        i,j = locarray.get_ij_from_xy(x,y)
                        curpar = 'Indices'
                        self.mytooltip.myparams[curgroup][curpar]={}
                        self.mytooltip.myparams[curgroup][curpar]['name'] ='Indice (i;j)'
                        self.mytooltip.myparams[curgroup][curpar]['value'] ='('+str(i+1)+';'+str(j+1)+')'
                        self.mytooltip.myparams[curgroup][curpar]['type'] ='String'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] =''

                    curpar = 'Value'
                    self.mytooltip.myparams[curgroup][curpar]={}
                    self.mytooltip.myparams[curgroup][curpar]['name'] ='Value'
                    self.mytooltip.myparams[curgroup][curpar]['value'] =val
                    self.mytooltip.myparams[curgroup][curpar]['type'] ='Float'
                    self.mytooltip.myparams[curgroup][curpar]['comment'] =''
                except:
                    pass
                
        if self.linked:
            for curFrame in self.linkedList:
                if not curFrame is self:
                    for locarray in curFrame.myarrays:
                        curgroup = locarray.idx
                        if curFrame.added['arrays'][curgroup]['checked']:
                            self.mytooltip.myparams[curgroup]={}

                            try:
                                val = locarray.get_value(x,y)
                                i,j = locarray.get_ij_from_xy(x,y)

                                curpar = 'Indices'
                                self.mytooltip.myparams[curgroup][curpar]={}
                                self.mytooltip.myparams[curgroup][curpar]['name'] ='Indice (i;j)'
                                self.mytooltip.myparams[curgroup][curpar]['value'] ='('+str(i+1)+';'+str(j+1)+')'
                                self.mytooltip.myparams[curgroup][curpar]['type'] ='String'
                                self.mytooltip.myparams[curgroup][curpar]['comment'] =''

                                curpar = 'Value'
                                self.mytooltip.myparams[curgroup][curpar]={}
                                self.mytooltip.myparams[curgroup][curpar]['name'] ='Value'
                                self.mytooltip.myparams[curgroup][curpar]['value'] =val
                                self.mytooltip.myparams[curgroup][curpar]['type'] ='Float'
                                self.mytooltip.myparams[curgroup][curpar]['comment'] =''
                            except:
                                pass

        self.mytooltip.PopulateOnePage()

        if  e.ControlDown():
            ttsize = self.mytooltip.GetSize()
            self.mytooltip.position(pos+posframe+(ttsize[0]/2.+15,15))
        else:
            width, height = self.GetSize()
            posframe[0]+=width
            self.mytooltip.position(posframe)

        self.mytooltip.Show(True)

        self.Refresh()

    def Autoscale(self,update_backfore=True):
        self.findminmax()
        self.width = self.xmax-self.xmin
        self.height = self.ymax-self.ymin

        centerx=self.xmin+self.width/2.
        centery=self.ymin+self.height/2.

        iwidth=self.width*self.sx
        iheight=self.height*self.sy

        width,height = self.canvas.GetSize()

        sx=float(width)/float(iwidth)
        sy=float(height)/float(iheight)

        if sx>sy:
            self.xmax=self.xmin + self.width*sx/sy
            self.width=self.xmax-self.xmin
        else:
            self.ymax=self.ymin + self.height*sy/sx
            self.height=self.ymax-self.ymin

        self.mousex=centerx
        self.mousey=centery

        if update_backfore:
            #dessin du background
            for id,(key,curmap) in enumerate(self.added['wms-background'].items()):
                if curmap['checked']:
                    curmap['values'].reload()
            #dessin du foreground
            for id,(key,curmap) in enumerate(self.added['wms-foreground'].items()):
                if curmap['checked']:
                    curmap['values'].reload()

        self.setbounds()
    
    def endactions(self):
        if self.action=='select by tmp vector' or self.action=='select by vector':            
            self.action=None
            self.active_vector.nbvertices=self.active_vector.nbvertices-1
            self.active_vector.myvertices.pop(-1)
            self.active_vector.close_force()
            self.active_array.mngselection.select_insidepoly(self.active_vector)
            
        if self.action=='capture vertices':
            self.action=None
            self.active_vector.nbvertices=self.active_vector.nbvertices-1
            self.active_vector.myvertices.pop(-1)
            r = wx.MessageDialog(
                                    None,
                                    _('End of points capturing') + '\n' +
                                    _('Force to close the vector ?'),
                                    _('Confirm'),
                                    wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION
                                ).ShowModal()
            if r ==wx.ID_YES:
                self.active_vector.close_force()
                
        if self.action=='modify vertices':
            self.active_vertex=None
            self.action=None
        
        if self.action=='insert vertices':
            self.active_vertex = None
            self.action = None
            
        if self.action=='dynamic parallel':
            self.active_vector.nbvertices-=1
            self.active_vector.myvertices.pop(-1)
            self.active_zone.parallel_active(self.dynapar_dist)
            
            self.active_zones.fill_structure()
            
            self.active_vertex = None
            self.action = None
            
        if self.action =='select active vector' or self.action =='select active vector2' or self.action=='select node by node':
            self.action = None
        
        self.copyfrom=None

        self.mimicme()
        
    def OnHotKey(self, e:wx.KeyEvent):
        '''
        Gestion des touches clavier
            F2 : mise à jour du résultat pas suivant
            F5 : autoscale
            F7 : refresh
            Z  : zoom avant
            z  : zoom artrière
            Flèches : déplacements latéraux
            P : sélection de profil
            
            !! ACTIONs !!
            N : sélection noeud par noeud de la matrice courante
            B : sélection par vecteur temporaire de la matrice courante
            V : sélection par vecteur activé de la matrice courante
            RETURN : fin d'action (cf aussi double clicks droit 'OnRDClick')
        '''
        key=e.GetKeyCode()
        
        if key==wx.WXK_CONTROL_U:
            msg=''
            if self.active_array is None:
                msg+=_('Active array is None\n')
                
            if msg!='':
                msg+=_('\n')
                msg+=_('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_fn_fnpos_gltf()                                                   
            self.update_blender_sculpting()
        
        elif key==wx.WXK_CONTROL_C:
            if self.active_array is None:
                dlg=wx.MessageDialog(self,_('The active array is None - Please active an array from which to copy the values !'),style=wx.OK)
                dlg.ShowModal()
                return

            self.copyfrom = self.active_array
            self.mimicme_copyfrom() # force le recopiage de copyfrom dans les autres matrices liées

        elif key == wx.WXK_CONTROL_V:
            if self.active_array is None:
                dlg=wx.MessageDialog(self,_('The active array is None - Please active an array into which to paste the values !'),style=wx.OK)
                dlg.ShowModal()
                return
            
            fromarray = self.copyfrom
            if fromarray is None:
                if self.linked:
                    if not self.linkedList is None:
                        for curFrame in self.linkedList:
                            if curFrame.copyfrom is not None:
                                fromarray = curFrame.copyfrom
                                break 
            if fromarray is None:
                dlg=wx.MessageDialog(self,_('No value to be pasted !'),style=wx.OK)
                dlg.ShowModal()
                return               
                
            cursel  = fromarray.mngselection.myselection
            if len(cursel)>0:
                z       = fromarray.mngselection.get_values_sel()
                self.active_array.set_values_sel(cursel,z)
            pass

        if key==wx.WXK_RETURN:
            self.endactions()
        
        if key==wx.WXK_F2:
            #Dessin des matrices
            try:
                for curarray in self.added['arrays']:
                    if self.added['arrays'][curarray]['checked']:
                        if type(self.added['arrays'][curarray]['values']) is wolfres2DGPU:
                            self.added['arrays'][curarray]['values'].onnext(1)
                
                self.Paint()
            except:
                pass
        elif key==wx.WXK_F5:
            #Autoscale
            self.Autoscale()
        elif key==wx.WXK_F7:
            self.update()
        elif key==wx.WXK_F12:
            self.active_array.myops.Show()
            self.active_array.myops.array_ops.SetSelection(1)
            self.active_array.myops.Center()
        elif key==ord('N'): #N
            if self.active_array is not None:
                self.active_array.myops.select_nod_by_node()
        elif key==ord('V'): #V
            if self.active_array is not None:
                self.active_array.myops.select_vector_inside_manager()
        elif key==ord('B'): #B
            if self.active_array is not None:
                self.active_array.myops.select_vector_tmp()            
        elif key ==ord('P'): #P
            self.action='Select nearest profile'
        elif key == ord('Z'): #Z
            self.width=self.width / 1.1
            self.height=self.height / 1.1
            self.setbounds()
        elif key == ord('z'): #z
            self.width=self.width * 1.1
            self.height=self.height * 1.1
            self.setbounds()
        elif key == wx.WXK_UP:
            self.mousey = self.mousey+self.height/10.
            self.setbounds()
        elif key == wx.WXK_DOWN:
            self.mousey = self.mousey-self.height/10.
            self.setbounds()
        elif key == wx.WXK_LEFT:
            self.mousex = self.mousex-self.width/10.
            self.setbounds()
        elif key == wx.WXK_RIGHT:
            self.mousex = self.mousex+self.width/10.
            self.setbounds()

    def update(self):
        #dessin du background
        for id,(key,curmap) in enumerate(self.added['wms-background'].items()):
            if curmap['checked']:
                curmap['values'].reload()
        #dessin du foreground
        for id,(key,curmap) in enumerate(self.added['wms-foreground'].items()):
            if curmap['checked']:
                curmap['values'].reload()
                
        if self.locminmax.IsChecked() or self.update_absolute_minmax:
            for id,(key,curmap) in enumerate(self.added['arrays'].items()):
                if curmap['checked']:
                    curarray:WolfArray
                    curarray=curmap['values']
                    if self.update_absolute_minmax:
                        curarray.updatepalette()
                        self.update_absolute_minmax=False
                    else:
                        curarray.updatepalette(onzoom=[self.xmin,self.xmax,self.ymin,self.ymax])
                    curarray.delete_lists()
        self.Refresh()

    def Paint(self):
        
        if self.currently_readresults:
            return

        width,height=self.canvas.GetSize()

        #C'est bien ici que la zone de dessin utile est calculée sur base du centre et de la zone en coordonnées réelles
        # Les commandes OpenGL sont donc traitées en coordonnées réelles puisque la commande glOrtho définit le cadre visible
        self.xmin=self.mousex-self.width/2.
        self.ymin=self.mousey-self.height/2.
        self.xmax=self.mousex+self.width/2.
        self.ymax=self.mousey+self.height/2.

        if self.canvas.SetCurrent(self.context):

            glClearColor(1.,1.,1.,0)
            glClear(GL_COLOR_BUFFER_BIT | GL_DEPTH_BUFFER_BIT)
            glViewport(0,0,int(width),int(height))

            glMatrixMode(GL_PROJECTION)
            glLoadIdentity()
            glOrtho(self.xmin, self.xmax, self.ymin, self.ymax, -99999, 99999)

            glMatrixMode(GL_MODELVIEW)
            glLoadIdentity()

            #dessin du background
            for id,(key,curmap) in enumerate(self.added['wms-background'].items()):
                if curmap['checked']:
                    curmap['values'].paint()

            #Dessin des matrices
            try:
                for curarray in self.added['arrays']:
                    if self.added['arrays'][curarray]['checked']:
                        locarray=self.added['arrays'][curarray]['values']
                        if not locarray.plotting:
                            locarray.plotting=True
                            self.added['arrays'][curarray]['values'].plot(self.sx,self.sy,self.xmin,self.ymin,self.xmax,self.ymax)
                            locarray.plotting=False
            except:
                pass

            #Dessin des résultats 2D
            try:
                for curarray in self.added['wolf2d']:
                    if self.added['wolf2d'][curarray]['checked']:
                        locarray=self.added['wolf2d'][curarray]['values']
                        if not locarray.plotting:
                            locarray.plotting=True
                            self.added['wolf2d'][curarray]['values'].plot(self.sx,self.sy,self.xmin,self.ymin,self.xmax,self.ymax)
                            locarray.plotting=False
            except:
                pass

            #Dessin des vecteurs
            try:
                for curvect in self.added['vectors']:
                    if self.added['vectors'][curvect]['checked']:
                        self.added['vectors'][curvect]['values'].plot()
            except:
                pass

            #Dessin des nuages
            try:
                for curcloud in self.added['clouds']:
                    if self.added['clouds'][curcloud]['checked']:
                        self.added['clouds'][curcloud]['values'].plot()
            except:
                pass

            #Dessin du reste
            try:
                for curobj in self.added['others']:
                    if self.added['others'][curobj]['checked']:
                        curobj=self.added['others'][curobj]['values']
                        if type(curobj) is SPWDCENNGaugingStations or type(curobj) is SPWMIGaugingStations:
                            curobj.plot((self.xmax-self.xmin)/100.)
                        elif type(curobj) is genericImagetexture:
                            curobj.paint()
                            
            except: 
                pass

            #Dessin du Front
            for id,(key,curmap) in enumerate(self.added['wms-foreground'].items()):
                if curmap['checked']:
                    curmap['values'].paint()


            #Gestion des BC (si actif)
            try:
                if self.mybc!=None:
                    self.mybc.plot()
            except:
                pass
        
            glFlush()
            self.canvas.SwapBuffers()
        else:
            raise NameError('Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def OnPaint(self, e):

        self.Paint()
        if e != None:
            e.Skip()

    def findminmax(self,force=False):
        xmin = 1.e30
        ymin = 1.e30
        xmax = -1.e30
        ymax = -1.e30

        k=0
        for locarray in self.myarrays:
            if self.added['arrays'][locarray.idx]['checked'] or force:
                xmin = min(locarray.origx+locarray.translx,xmin)
                xmax = max(locarray.origx+locarray.translx + float(locarray.nbx) * locarray.dx,xmax)
                ymin = min(locarray.origy+locarray.transly,ymin)
                ymax = max(locarray.origy+locarray.transly + float(locarray.nby) * locarray.dy,ymax)
                k+=1

        for locvector in self.myvectors:
            if self.added['vectors'][locvector.idx]['checked'] or force:
                if locvector.idx!='grid':
                    locvector.find_minmax()
                    xmin = min(locvector.minx,xmin)
                    xmax = max(locvector.maxx,xmax)
                    ymin = min(locvector.miny,ymin)
                    ymax = max(locvector.maxy,ymax)
                    k+=1

        for loccloud in self.myclouds:
            if self.added['clouds'][loccloud.idx]['checked'] or force:
                loccloud.find_minmax(force)
                xmin = min(loccloud.xbounds[0],xmin)
                xmax = max(loccloud.xbounds[1],xmax)
                ymin = min(loccloud.ybounds[0],ymin)
                ymax = max(loccloud.ybounds[1],ymax)
                k+=1

        for locres2d in self.myres2D:
            if self.added['wolf2d'][locres2d.idx]['checked'] or force:
                xmin = min(locres2d.origx+locres2d.translx,xmin)
                xmax = max(locres2d.origx+locres2d.translx + float(locres2d.nbx) * locres2d.dx,xmax)
                ymin = min(locres2d.origy+locres2d.transly,ymin)
                ymax = max(locres2d.origy+locres2d.transly + float(locres2d.nby) * locres2d.dy,ymax)
                k+=1

        for locothers in self.myothers:
            if type(locothers) is genericImagetexture:
                xmin = locothers.xmin
                xmax = locothers.xmax
                ymin = locothers.ymin
                ymax = locothers.ymax
                k+=1

        if k>0:
            self.xmin=xmin
            self.xmax=xmax
            self.ymin=ymin
            self.ymax=ymax
    
    def resizeFrame(self,w,h):
        self.SetClientSize(w,h)
        
    def mimicme(self):
        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width,height=self.GetClientSize()
                
                curFrame:WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic=False
                
                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.resizeFrame(width,height)
                        curFrame.mousex = self.mousex
                        curFrame.mousey = self.mousey
                        curFrame.sx = self.sx
                        curFrame.sy = self.sy
                        curFrame.width = self.width
                        curFrame.height = self.height
                        curFrame.setbounds()
                        
                        if curFrame.link_shareopsvect:
                            curFrame.Active_vector(self.active_vector)
                            curFrame.active_array.myops.Active_vector(self.active_vector,False)
                            curFrame.action=self.action

                for curFrame in self.linkedList:
                    curFrame.forcemimic=True

    def mimicme_copyfrom(self):
        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width,height=self.GetClientSize()
                
                curFrame:WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic=False
                
                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.copyfrom = self.copyfrom

                for curFrame in self.linkedList:
                    curFrame.forcemimic=True

    def Active_vector(self,vect):
        self.active_vector = vect    
        if vect is not None:
            if vect.parentzone is not None:
                self.Active_zone(vect.parentzone)
            
        self.mimicme()

    def Active_zone(self,zone:zone):
        self.active_zone = zone
        self.active_zones = zone.parent

class genericImagetexture():
    
    parent=None

    name:str
    idtexture:int
    xmin:float
    xmax:float
    ymin:float
    ymax:float

    width:int
    height:int

    which:str

    myImage:Image

    def __init__(self,which:str,label:str,parent,xmin,xmax,ymin,ymax,width=1000,height=1000, imageFile="", imageObj=None) -> None:
        
        parent.canvas.SetCurrent(parent.context)
        self.parent=parent

        self.xmin=xmin
        self.xmax=xmax
        self.ymin=ymin
        self.ymax=ymax
        self.idtexture=(GLuint*1)()
        try:
            glGenTextures(1,self.idtexture)
        except:
            raise NameError('Opengl glGenTextures -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')
        self.width=width
        self.height=height
        self.which=which.lower()
        self.name=label
        self.oldview=[self.xmin,self.xmax,self.ymin,self.ymax,self.width,self.height]
        self.imageFile = imageFile
        # TO CHECK !!!!
        self.myImage = imageObj

        self.load(self.imageFile)        

        pass

    def load(self, imageFile=""):
        if self.width==-99999 or self.height==-99999:
            return

        if self.parent.canvas.SetCurrent(self.parent.context):
            mybytes:BytesIO

            # myFile = "F:\\Christophe\\Simulations\\Optimisation\\Amblève\\Optimisation_auto\\SA\\Reference\\Test.png"
            if imageFile!="":
                self.myImage=Image.open(imageFile)

            glBindTexture(GL_TEXTURE_2D, self.idtexture[0])
            glTexImage2D(GL_TEXTURE_2D, 0, GL_RGBA, self.myImage.width,self.myImage.height, 0, GL_RGBA, GL_UNSIGNED_BYTE, self.myImage.tobytes())
            
            glTexParameteri(GL_TEXTURE_2D, GL_TEXTURE_MAG_FILTER, GL_LINEAR)
            glTexParameteri(GL_TEXTURE_2D, GL_TEXTURE_MIN_FILTER, GL_LINEAR_MIPMAP_LINEAR)
            glGenerateMipmap(GL_TEXTURE_2D)
        else:
            raise NameError('Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')
            
    def reload(self):
        dx=self.parent.xmax-self.parent.xmin
        dy=self.parent.ymax-self.parent.ymin
        try:
            if dx!=self.myImage.width:
                dx = max(math.floor(dx/self.myImage.width),1)*self.myImage.width

            if dy!=self.myImage.height:
                dy = max(math.floor(dy/self.myImage.height),1)*self.myImage.height
        except:
            pass
            

        cx=self.parent.mousex
        cy=self.parent.mousey

        coeff=.5
        # self.xmin=cx-dx*coeff
        # self.xmax=cx+dx*coeff
        # self.ymin=cy-dy*coeff
        # self.ymax=cy+dy*coeff
        # self.width=self.parent.canvaswidth*2*coeff
        # self.height=self.parent.canvasheight*2*coeff

        self.newview=[self.xmin,self.xmax,self.ymin,self.ymax,self.width,self.height]
        if self.newview!=self.oldview:
            self.load()        
            self.oldview=self.newview

    def paint(self):
        
        glColor4f(1.,1.,1.,1.)
        glEnable(GL_TEXTURE_2D)
        glEnable(GL_BLEND)
        glBlendFunc(GL_SRC_ALPHA, GL_ONE_MINUS_SRC_ALPHA)

        glBindTexture(GL_TEXTURE_2D, self.idtexture[0])

        glBegin(GL_QUADS)
        glTexCoord2f(0.0, 0.0)
        glVertex2f(self.xmin,self.ymax)
        glTexCoord2f(1.0, 0.0)
        glVertex2f(self.xmax,self.ymax)
        glTexCoord2f(1.0, 1.0)
        glVertex2f(self.xmax,self.ymin)
        glTexCoord2f(0.0, 1.0)
        glVertex2f(self.xmin,self.ymin)
        glEnd()
        glDisable(GL_TEXTURE_2D)
        glDisable(GL_BLEND)
