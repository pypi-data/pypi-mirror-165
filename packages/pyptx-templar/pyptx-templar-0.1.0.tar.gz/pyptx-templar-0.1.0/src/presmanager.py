import logging
import os
import sys

import copy

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.table import _Column, _Row

from pyptx_templar.copy import textframe_copy


def export(pres, outpath):
    try:
        os.remove(outpath)
        logging.info(f"Former {outpath} deleted.")
    except OSError:
        pass
    pres.save(outpath)


def remove_slide(pres, idx):
    rId = pres.slides._sldIdLst[idx].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[idx]


def delete_run(r):
    r._r.getparent().remove(r._r)


def clear_slides(pres):
    for idx in range(len(pres.slides) - 1, -1, -1):
        remove_slide(pres, idx)


def insert_slide(pres, layout, idx):
    slide = append_slide(pres, layout)
    move_slide(pres.slides, slide, idx)
    return slide


def move_slide(slides, slide, new_idx):
    slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])


def fill_placeholders(slide):
    layout = slide.slide_layout
    for ph in layout.placeholders:
        if ph.has_text_frame:
            try:
                phid = ph.placeholder_format.idx
                tf = slide.placeholders[phid].text_frame
            except KeyError:
                continue
            textframe_copy(ph.text_frame, tf)


# Ajoute la slide avec le nom donné à la fin de la présentation
def append_slide(pres, layout):
    if isinstance(layout, str):
        nl = pres.slide_masters[0].slide_layouts.get_by_name(layout)
        if nl is None:
            logging.error("append_slide: layout \'%s\' inexistant", layout)
            sys.exit(1)
        layout = nl

    slide = pres.slides.add_slide(layout)
    fill_placeholders(slide)
    return slide


# Ecrit le numéro des shapes, pratique pour repérer initialement
def etalonnage(prespath, outpath):
    pres = Presentation(prespath)
    for idx, slide in enumerate(pres.slides):
        logging.info("Slide %d has slide_id %s", idx, slide.slide_id)

    clear_slides(pres)
    for layout in pres.slide_masters[0].slide_layouts:
        slide = pres.slides.add_slide(layout)
        for ph in slide.placeholders:
            ph.text = f"placeholder {ph.placeholder_format.idx}"
    export(pres, outpath)


def chart_replace_data(chart, values):
    chart_data = CategoryChartData()
    chart_data.categories = chart.plots[0].categories
    chart_data.add_series('Values', values)
    chart.replace_data(chart_data)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


# duplique la dernière ligne, en particulier les bordures et le texte
# sont gardés exactement à l'identique
def table_dup_row(table, idx, n=1, to=None):
    if to is None:
        to = idx
    to = (to + len(table.rows)) % len(table.rows)

    new_row = table._tbl.tr_lst[idx]
    num = -1
    for rid, rxml in enumerate(table._tbl.iterchildren()):
        if rxml.tag == new_row.tag:
            num += 1
            if num == to:
                break

    for _ in range(n):
        new_row = copy.deepcopy(new_row)
        table._tbl.insert(rid, new_row)

    if n == 1:
        return _Row(new_row, table)
    return None


def table_dup_column(table, idx, n=1, to=None, keep_width=False):
    total_width = table._graphic_frame.width
    colwidth = table.columns[idx].width
    mult = total_width / (total_width + n * colwidth)

    if to is None:
        to = idx
    to = (to + len(table.columns)) % len(table.columns)

    new_col = table._tbl.tblGrid.gridCol_lst[idx]
    num = -1
    for cid, cxml in enumerate(table._tbl.tblGrid.iterchildren()):
        if cxml.tag == new_col.tag:
            num += 1
            if num == to:
                break
    for _ in range(n):
        new_col = copy.deepcopy(new_col)
        table._tbl.tblGrid.insert(cid, new_col)

    for tr in table._tbl.tr_lst:
        new_cell = tr.tc_lst[idx]
        num = -1
        for cid, cxml in enumerate(tr.iterchildren()):
            if cxml.tag == new_cell.tag:
                num += 1
                if num == to:
                    break
        for _ in range(n):
            new_cell = copy.deepcopy(new_cell)
            tr.insert(cid, new_cell)

    if keep_width:
        for col in table.columns:
            col.width = int(col.width * mult)

    if n == 1:
        return _Column(new_col, table)
    return None


# def _get_blank_slide_layout(pres):
#     layout_items_count = [
#         len(layout.placeholders) for layout in pres.slide_layouts
#     ]
#     min_items = min(layout_items_count)
#     blank_layout_id = layout_items_count.index(min_items)
#     return pres.slide_layouts[blank_layout_id]


def duplicate_slide(pres, index):
    source = pres.slides[index]
    dest = pres.slides.add_slide(source.slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in source.rels.items():
        if not "notesSlide" in value.reltype:
            dest.rels.add_relationship(value.reltype, value._target, value.rId)

    move_slide(pres.slides, pres.slides[-1], index)

    return dest
